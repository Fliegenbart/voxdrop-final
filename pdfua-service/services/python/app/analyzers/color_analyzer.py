"""
PPTX Color Analyzer - Extracts all colors from PowerPoint presentations.
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE, MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from typing import List, Dict, Optional, Any, Tuple
import logging

logger = logging.getLogger(__name__)


def rgb_to_hex(r: int, g: int, b: int) -> str:
    """Convert RGB values to hex string."""
    return f"#{r:02x}{g:02x}{b:02x}"


def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """Convert hex string to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


class ColorExtractor:
    """Extracts all colors from a PowerPoint presentation."""

    # Default theme color approximations (when theme can't be resolved)
    DEFAULT_THEME_COLORS = {
        MSO_THEME_COLOR.DARK_1: (0, 0, 0),         # Usually black
        MSO_THEME_COLOR.LIGHT_1: (255, 255, 255),  # Usually white
        MSO_THEME_COLOR.DARK_2: (68, 84, 106),     # Dark gray-blue
        MSO_THEME_COLOR.LIGHT_2: (231, 230, 230),  # Light gray
        MSO_THEME_COLOR.ACCENT_1: (68, 114, 196),  # Blue
        MSO_THEME_COLOR.ACCENT_2: (237, 125, 49),  # Orange
        MSO_THEME_COLOR.ACCENT_3: (165, 165, 165), # Gray
        MSO_THEME_COLOR.ACCENT_4: (255, 192, 0),   # Yellow
        MSO_THEME_COLOR.ACCENT_5: (91, 155, 213),  # Light blue
        MSO_THEME_COLOR.ACCENT_6: (112, 173, 71),  # Green
    }

    def extract_all_colors(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract all colors from a PowerPoint presentation.

        Returns:
            Dict with slides data including colors and their locations.
        """
        try:
            pres = Presentation(pptx_path)
        except Exception as e:
            logger.error(f"Failed to open PPTX: {e}")
            raise ValueError(f"Could not open PowerPoint file: {e}")

        slides_data = []
        all_color_pairs = []

        for slide_num, slide in enumerate(pres.slides, 1):
            slide_data = {
                'slide_number': slide_num,
                'background': self._get_background(slide),
                'text_elements': self._get_text_elements(slide),
                'shape_fills': self._get_shape_fills(slide),
                'color_pairs': [],
                'issues': []
            }

            # Analyze color pairs on this slide
            bg_color = slide_data['background']

            for text_elem in slide_data['text_elements']:
                fg_color = text_elem['color']
                pair = {
                    'foreground': fg_color,
                    'background': bg_color,
                    'text_preview': text_elem.get('text', ''),
                    'location': f"Folie {slide_num}",
                    'type': 'text'
                }
                slide_data['color_pairs'].append(pair)
                all_color_pairs.append(pair)

            # Check shape fills against background
            for shape_fill in slide_data['shape_fills']:
                if shape_fill.get('has_text'):
                    # Shape with text - check text color against fill
                    for text_color in shape_fill.get('text_colors', []):
                        pair = {
                            'foreground': text_color,
                            'background': shape_fill['color'],
                            'text_preview': shape_fill.get('text_preview', ''),
                            'location': f"Folie {slide_num} (Shape)",
                            'type': 'shape_text'
                        }
                        slide_data['color_pairs'].append(pair)
                        all_color_pairs.append(pair)

            slides_data.append(slide_data)

        return {
            'slides': slides_data,
            'total_slides': len(slides_data),
            'all_color_pairs': all_color_pairs
        }

    def _get_background(self, slide) -> Dict[str, Any]:
        """Extract slide background color."""
        default_bg = {'r': 255, 'g': 255, 'b': 255, 'hex': '#ffffff'}

        try:
            fill = slide.background.fill

            if fill.type == MSO_FILL_TYPE.SOLID:
                color_format = fill.fore_color

                # Try to get RGB directly
                if hasattr(color_format, 'rgb') and color_format.rgb:
                    rgb = color_format.rgb
                    return {
                        'r': rgb.red if hasattr(rgb, 'red') else rgb[0],
                        'g': rgb.green if hasattr(rgb, 'green') else rgb[1],
                        'b': rgb.blue if hasattr(rgb, 'blue') else rgb[2],
                        'hex': rgb_to_hex(
                            rgb.red if hasattr(rgb, 'red') else rgb[0],
                            rgb.green if hasattr(rgb, 'green') else rgb[1],
                            rgb.blue if hasattr(rgb, 'blue') else rgb[2]
                        )
                    }

                # Try theme color
                if hasattr(color_format, 'theme_color') and color_format.theme_color:
                    theme_rgb = self.DEFAULT_THEME_COLORS.get(
                        color_format.theme_color,
                        (255, 255, 255)
                    )
                    return {
                        'r': theme_rgb[0],
                        'g': theme_rgb[1],
                        'b': theme_rgb[2],
                        'hex': rgb_to_hex(*theme_rgb),
                        'theme': True
                    }
        except Exception as e:
            logger.debug(f"Could not extract background color: {e}")

        return default_bg

    def _get_text_elements(self, slide) -> List[Dict[str, Any]]:
        """Extract all text elements with their colors."""
        elements = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue

                    color = self._extract_run_color(run)
                    if color:
                        elements.append({
                            'text': run.text[:100],
                            'color': color,
                            'font_size': run.font.size.pt if run.font.size else None
                        })

        return elements

    def _extract_run_color(self, run) -> Optional[Dict[str, Any]]:
        """Extract color from a text run."""
        default_color = {'r': 0, 'g': 0, 'b': 0, 'hex': '#000000'}

        try:
            color_format = run.font.color

            # Direct RGB
            if hasattr(color_format, 'rgb') and color_format.rgb:
                rgb = color_format.rgb
                return {
                    'r': rgb.red if hasattr(rgb, 'red') else rgb[0],
                    'g': rgb.green if hasattr(rgb, 'green') else rgb[1],
                    'b': rgb.blue if hasattr(rgb, 'blue') else rgb[2],
                    'hex': rgb_to_hex(
                        rgb.red if hasattr(rgb, 'red') else rgb[0],
                        rgb.green if hasattr(rgb, 'green') else rgb[1],
                        rgb.blue if hasattr(rgb, 'blue') else rgb[2]
                    )
                }

            # Theme color
            if hasattr(color_format, 'theme_color') and color_format.theme_color:
                theme_rgb = self.DEFAULT_THEME_COLORS.get(
                    color_format.theme_color,
                    (0, 0, 0)
                )
                return {
                    'r': theme_rgb[0],
                    'g': theme_rgb[1],
                    'b': theme_rgb[2],
                    'hex': rgb_to_hex(*theme_rgb),
                    'theme': True
                }
        except Exception as e:
            logger.debug(f"Could not extract text color: {e}")

        return default_color

    def _get_shape_fills(self, slide) -> List[Dict[str, Any]]:
        """Extract shape fill colors."""
        fills = []

        for shape in slide.shapes:
            try:
                if not hasattr(shape, 'fill'):
                    continue

                fill = shape.fill
                if fill.type != MSO_FILL_TYPE.SOLID:
                    continue

                color_format = fill.fore_color
                shape_data = {
                    'has_text': shape.has_text_frame and bool(shape.text_frame.text.strip()),
                    'text_colors': [],
                    'text_preview': ''
                }

                # Get fill color
                if hasattr(color_format, 'rgb') and color_format.rgb:
                    rgb = color_format.rgb
                    shape_data['color'] = {
                        'r': rgb.red if hasattr(rgb, 'red') else rgb[0],
                        'g': rgb.green if hasattr(rgb, 'green') else rgb[1],
                        'b': rgb.blue if hasattr(rgb, 'blue') else rgb[2],
                        'hex': rgb_to_hex(
                            rgb.red if hasattr(rgb, 'red') else rgb[0],
                            rgb.green if hasattr(rgb, 'green') else rgb[1],
                            rgb.blue if hasattr(rgb, 'blue') else rgb[2]
                        )
                    }
                elif hasattr(color_format, 'theme_color') and color_format.theme_color:
                    theme_rgb = self.DEFAULT_THEME_COLORS.get(
                        color_format.theme_color,
                        (200, 200, 200)
                    )
                    shape_data['color'] = {
                        'r': theme_rgb[0],
                        'g': theme_rgb[1],
                        'b': theme_rgb[2],
                        'hex': rgb_to_hex(*theme_rgb),
                        'theme': True
                    }
                else:
                    continue

                # Get text colors if shape has text
                if shape_data['has_text']:
                    shape_data['text_preview'] = shape.text_frame.text[:50]
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            color = self._extract_run_color(run)
                            if color and color not in shape_data['text_colors']:
                                shape_data['text_colors'].append(color)

                fills.append(shape_data)

            except Exception as e:
                logger.debug(f"Could not extract shape fill: {e}")
                continue

        return fills
