"""
PPTX Color Fixer - Replaces colors in PowerPoint presentations.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL_TYPE
from typing import Tuple, List, Dict, Any, Optional
import logging
import os

from .contrast import suggest_better_contrast, calculate_contrast_ratio, calculate_luminance

logger = logging.getLogger(__name__)


class ColorFixer:
    """Fixes colors in PowerPoint presentations for better accessibility."""

    def __init__(self, similarity_threshold: int = 15):
        """
        Initialize ColorFixer.

        Args:
            similarity_threshold: Max RGB difference to consider colors as matching
        """
        self.similarity_threshold = similarity_threshold

    def apply_fixes(
        self,
        pptx_path: str,
        replacements: List[Dict[str, Any]],
        output_path: Optional[str] = None
    ) -> str:
        """
        Apply color replacements to a PowerPoint file.

        Args:
            pptx_path: Path to input PPTX file
            replacements: List of replacement dicts with:
                - old: [r, g, b] - Color to replace
                - new: [r, g, b] - New color
                - target: 'all', 'text', 'shapes', 'background'
                - slide: Optional slide number (None = all slides)
            output_path: Path for output file (defaults to input_fixed.pptx)

        Returns:
            Path to the modified PPTX file
        """
        try:
            pres = Presentation(pptx_path)
        except Exception as e:
            logger.error(f"Failed to open PPTX: {e}")
            raise ValueError(f"Could not open PowerPoint file: {e}")

        for repl in replacements:
            old_rgb = tuple(repl['old'])
            new_rgb = RGBColor(*repl['new'])
            target = repl.get('target', 'all')
            slide_num = repl.get('slide')

            for idx, slide in enumerate(pres.slides, 1):
                # Skip if specific slide requested and this isn't it
                if slide_num is not None and idx != slide_num:
                    continue

                if target in ['all', 'text']:
                    self._fix_text_colors(slide, old_rgb, new_rgb)

                if target in ['all', 'shapes']:
                    self._fix_shape_fills(slide, old_rgb, new_rgb)

                if target in ['all', 'background']:
                    self._fix_background(slide, old_rgb, new_rgb)

        # Determine output path
        if output_path is None:
            base, ext = os.path.splitext(pptx_path)
            output_path = f"{base}_fixed{ext}"

        pres.save(output_path)
        logger.info(f"Saved fixed PPTX to {output_path}")

        return output_path

    def auto_fix_contrast(
        self,
        pptx_path: str,
        target_level: str = 'aa',
        output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Automatically fix all contrast issues in a PowerPoint file.

        Args:
            pptx_path: Path to input PPTX file
            target_level: 'aa' (4.5:1) or 'aaa' (7:1)
            output_path: Path for output file

        Returns:
            Dict with fix report and output path
        """
        target_ratio = 7.0 if target_level == 'aaa' else 4.5

        try:
            pres = Presentation(pptx_path)
        except Exception as e:
            raise ValueError(f"Could not open PowerPoint file: {e}")

        fixes_applied = []

        for slide_idx, slide in enumerate(pres.slides, 1):
            # Get slide background color
            bg_color = self._get_slide_background(slide)

            # Fix text colors
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # Check if shape has its own fill
                shape_bg = self._get_shape_fill(shape)
                bg_was_detected = shape_bg is not None or bg_color != (255, 255, 255)
                effective_bg = shape_bg if shape_bg else bg_color

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text.strip():
                            continue

                        fg_color = self._get_run_color(run)
                        if fg_color is None:
                            continue

                        # Safety check: If text is very light (like white) and we couldn't
                        # detect a background, assume the background is dark and skip.
                        # This prevents turning white-on-dark-blue text into black-on-dark-blue.
                        fg_luminance = calculate_luminance(fg_color)
                        if not bg_was_detected and fg_luminance > 0.7:
                            logger.debug(
                                f"Skipping light text (lum={fg_luminance:.2f}) on undetected background"
                            )
                            continue

                        current_ratio = calculate_contrast_ratio(fg_color, effective_bg)

                        if current_ratio < target_ratio:
                            suggestion = suggest_better_contrast(
                                fg_color, effective_bg, target_ratio
                            )

                            if suggestion['changed'] and suggestion['meets_target']:
                                new_color = suggestion['color']
                                run.font.color.rgb = RGBColor(
                                    new_color['r'],
                                    new_color['g'],
                                    new_color['b']
                                )

                                fixes_applied.append({
                                    'slide': slide_idx,
                                    'type': 'text',
                                    'text': run.text[:30],
                                    'old_color': f"#{fg_color[0]:02x}{fg_color[1]:02x}{fg_color[2]:02x}",
                                    'new_color': suggestion['hex'],
                                    'old_ratio': round(current_ratio, 2),
                                    'new_ratio': suggestion['ratio']
                                })

        # Save output
        if output_path is None:
            base, ext = os.path.splitext(pptx_path)
            output_path = f"{base}_fixed{ext}"

        pres.save(output_path)

        return {
            'output_path': output_path,
            'fixes_applied': fixes_applied,
            'total_fixes': len(fixes_applied),
            'target_level': target_level.upper()
        }

    def _fix_text_colors(
        self,
        slide,
        old_rgb: Tuple[int, int, int],
        new_rgb: RGBColor
    ):
        """Fix text colors in a slide."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    try:
                        current = self._get_run_color(run)
                        if current and self._colors_match(current, old_rgb):
                            run.font.color.rgb = new_rgb
                    except Exception as e:
                        logger.debug(f"Could not fix text color: {e}")

    def _fix_shape_fills(
        self,
        slide,
        old_rgb: Tuple[int, int, int],
        new_rgb: RGBColor
    ):
        """Fix shape fill colors in a slide."""
        for shape in slide.shapes:
            try:
                if not hasattr(shape, 'fill'):
                    continue

                if shape.fill.type != MSO_FILL_TYPE.SOLID:
                    continue

                current = self._get_shape_fill(shape)
                if current and self._colors_match(current, old_rgb):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = new_rgb
            except Exception as e:
                logger.debug(f"Could not fix shape fill: {e}")

    def _fix_background(
        self,
        slide,
        old_rgb: Tuple[int, int, int],
        new_rgb: RGBColor
    ):
        """Fix slide background color."""
        try:
            current = self._get_slide_background(slide)
            if current and self._colors_match(current, old_rgb):
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = new_rgb
        except Exception as e:
            logger.debug(f"Could not fix background: {e}")

    def _get_run_color(self, run) -> Optional[Tuple[int, int, int]]:
        """Extract color from a text run."""
        try:
            if run.font.color.rgb:
                rgb = run.font.color.rgb
                if hasattr(rgb, 'red'):
                    return (rgb.red, rgb.green, rgb.blue)
                return tuple(rgb)
        except:
            pass
        return None

    def _get_shape_fill(self, shape) -> Optional[Tuple[int, int, int]]:
        """Get fill color of a shape, supporting multiple fill types."""
        try:
            if not hasattr(shape, 'fill'):
                return None

            fill = shape.fill
            fill_type = fill.type

            # Handle solid fill
            if fill_type == MSO_FILL_TYPE.SOLID:
                rgb = fill.fore_color.rgb
                if rgb:
                    if hasattr(rgb, 'red'):
                        return (rgb.red, rgb.green, rgb.blue)
                    return tuple(rgb)
                # Try theme color
                try:
                    theme_color = fill.fore_color.theme_color
                    if theme_color is not None:
                        # Try to get the actual RGB value
                        brightness = getattr(fill.fore_color, 'brightness', 0) or 0
                        # Estimate color from theme (dark themes are common)
                        if brightness < 0:  # Darker
                            return (50, 50, 100)  # Assume dark blue-ish
                        elif brightness > 0.5:  # Lighter
                            return (200, 200, 200)
                except:
                    pass

            # Handle gradient fill - use the first gradient stop color
            elif fill_type == MSO_FILL_TYPE.GRADIENT:
                try:
                    stops = fill.gradient_stops
                    if stops and len(stops) > 0:
                        rgb = stops[0].color.rgb
                        if rgb:
                            if hasattr(rgb, 'red'):
                                return (rgb.red, rgb.green, rgb.blue)
                            return tuple(rgb)
                except:
                    pass

        except Exception as e:
            logger.debug(f"Could not get shape fill: {e}")

        return None

    def _get_slide_background(self, slide) -> Tuple[int, int, int]:
        """Get slide background color, defaulting to white."""
        try:
            fill = slide.background.fill
            fill_type = fill.type

            if fill_type == MSO_FILL_TYPE.SOLID:
                rgb = fill.fore_color.rgb
                if rgb:
                    if hasattr(rgb, 'red'):
                        return (rgb.red, rgb.green, rgb.blue)
                    return tuple(rgb)

            # Handle gradient fill
            elif fill_type == MSO_FILL_TYPE.GRADIENT:
                try:
                    stops = fill.gradient_stops
                    if stops and len(stops) > 0:
                        rgb = stops[0].color.rgb
                        if rgb:
                            if hasattr(rgb, 'red'):
                                return (rgb.red, rgb.green, rgb.blue)
                            return tuple(rgb)
                except:
                    pass

        except:
            pass
        return (255, 255, 255)  # Default white

    def _colors_match(
        self,
        c1: Tuple[int, int, int],
        c2: Tuple[int, int, int]
    ) -> bool:
        """Check if two colors are similar within threshold."""
        return all(
            abs(a - b) <= self.similarity_threshold
            for a, b in zip(c1, c2)
        )
