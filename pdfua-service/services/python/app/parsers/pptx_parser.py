"""
PPTX Parser - Extracts slides, text, and images from PowerPoint files.
"""
import io
import base64
import re
from typing import List, Dict, Any, Optional, Tuple, Set
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from PIL import Image


# Placeholder types that are considered footers (should be skipped)
FOOTER_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.SLIDE_NUMBER,
}

# Placeholder types considered as title anchors.
TITLE_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
}

PLACEHOLDER_RESIDUE_PATTERNS = (
    "klicken sie, um text hinzuzufügen",
    "klicken sie hier, um text hinzuzufügen",
    "klicken sie, um titel hinzuzufügen",
    "click to add text",
    "click to add title",
    "click to add subtitle",
)


def parse_pptx(
    file_path: str,
    *,
    include_images: bool = True,
    output_mode: str = "faithful_accessible",
    include_speaker_notes: bool = True,
) -> List[Dict[str, Any]]:
    """
    Parse a PPTX file and extract slide content.

    Returns a list of slide dictionaries with:
    - slide_number: int
    - title: str | None
    - text_content: list of text blocks
    - images: list of image data
    - shapes: list of shape info
    - speaker_notes: str | None
    """
    prs = Presentation(file_path)
    slides = []
    slide_width = float(prs.slide_width or 0)
    slide_height = float(prs.slide_height or 0)

    for slide_num, slide in enumerate(prs.slides, start=1):
        if _is_hidden_slide(slide):
            continue
        slide_data = {
            "slide_number": slide_num,
            "title": None,
            "title_source": None,
            "text_content": [],
            "images": [],
            "shapes": [],
            "speaker_notes": None,
            "speaker_notes_visibility": "context_only" if include_speaker_notes else "ignored",
            "has_table": False,
            "has_chart": False,
            "complexity": None,  # Will be set by classifier
            "slide_width": slide_width,
            "slide_height": slide_height,
            "output_mode": output_mode,
            "risk_flags": [],
            "suppressed_content": [],
        }

        # Extract title (primary + robust placeholder fallback).
        try:
            if slide.shapes.title:
                title_text = _normalize_pptx_text(slide.shapes.title.text)
                if title_text:
                    slide_data["title"] = title_text
                    slide_data["title_source"] = "shape_title"
        except Exception:
            pass

        if not slide_data.get("title"):
            fallback_title = _extract_slide_title_fallback(slide, slide_width, slide_height)
            if fallback_title:
                slide_data["title"] = fallback_title
                slide_data["title_source"] = "placeholder_fallback"

        # Process all shapes
        for shape in slide.shapes:
            shape_info = process_shape(
                shape,
                slide_num,
                slide_width,
                slide_height,
                include_images=include_images,
            )
            if shape_info:
                if shape_info["type"] == "suppressed":
                    risk_flag = str(shape_info.get("risk_flag") or "").strip()
                    if risk_flag and risk_flag not in slide_data["risk_flags"]:
                        slide_data["risk_flags"].append(risk_flag)
                    slide_data["suppressed_content"].append({
                        "reason": shape_info.get("reason"),
                        "risk_flag": risk_flag or None,
                        "preview": shape_info.get("preview"),
                        "visibility_source": shape_info.get("visibility_source"),
                    })
                elif shape_info["type"] == "text":
                    shape_info["index"] = len(slide_data["text_content"])
                    slide_data["text_content"].append(shape_info)
                elif shape_info["type"] == "image":
                    slide_data["images"].append(shape_info)
                elif shape_info["type"] == "table":
                    slide_data["has_table"] = True
                    slide_data["shapes"].append(shape_info)
                elif shape_info["type"] == "chart":
                    slide_data["has_chart"] = True
                    slide_data["shapes"].append(shape_info)
                else:
                    slide_data["shapes"].append(shape_info)

        # Extract speaker notes
        if include_speaker_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data["speaker_notes"] = notes_text

        # Derive a better reading order + structure hints for downstream LLM and HTML.
        try:
            ordered_blocks = sort_text_blocks_by_layout(slide_data["text_content"])
            slide_data["text_ordered"] = ordered_blocks
            structured_items, used_indices = detect_numbered_groups(ordered_blocks)
            if structured_items:
                slide_data["structured_items"] = structured_items
                slide_data["structured_text"] = build_structured_text(structured_items)
                slide_data["structured_used_indices"] = sorted(used_indices)
        except Exception as layout_err:
            print(f"[PPTX Parser] Layout analysis failed on slide {slide_num}: {layout_err}")

        slides.append(slide_data)

    return slides


def extract_pptx_metadata(file_path: str) -> Dict[str, Any]:
    """Extract document-level metadata from PPTX core properties."""
    try:
        prs = Presentation(file_path)
        props = prs.core_properties
        return {
            "title": props.title,
            "author": props.author,
            "last_modified_by": props.last_modified_by,
            "subject": props.subject,
            "category": props.category,
            "keywords": props.keywords,
            "comments": props.comments,
            "created": props.created.isoformat() if props.created else None,
            "modified": props.modified.isoformat() if props.modified else None,
            "revision": props.revision,
        }
    except Exception as e:
        print(f"[PPTX Parser] Metadata extraction failed: {e}")
        return {}


def is_footer_shape(shape) -> bool:
    """Check if a shape is a footer placeholder (date, slide number, footer text)."""
    try:
        # Check if it's a placeholder with footer type
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type in FOOTER_PLACEHOLDER_TYPES:
                return True

        # Also check by position - footers are usually at the very bottom
        # and have small height (less than 5% of slide height)
        if hasattr(shape, 'top') and hasattr(shape, 'height'):
            # Typical slide height is ~5143500 EMUs (7.5 inches)
            # Footer is usually in bottom 10%
            slide_height = 6858000  # ~9.5 inches in EMUs (standard height)
            shape_bottom = shape.top + shape.height
            if shape_bottom > slide_height * 0.9:
                # Small text at bottom - likely footer
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    # Common footer patterns
                    if len(text) < 50:  # Footer text is usually short
                        return True

        return False
    except Exception:
        return False


def _normalize(value: float, total: float) -> float:
    if not total:
        return 0.0
    return float(value) / float(total)


def _shape_position(shape, slide_width: float, slide_height: float) -> Dict[str, float]:
    """Extract raw + normalized bounds for layout-aware ordering."""
    left = float(getattr(shape, "left", 0) or 0)
    top = float(getattr(shape, "top", 0) or 0)
    width = float(getattr(shape, "width", 0) or 0)
    height = float(getattr(shape, "height", 0) or 0)

    left_norm = _normalize(left, slide_width)
    top_norm = _normalize(top, slide_height)
    width_norm = _normalize(width, slide_width)
    height_norm = _normalize(height, slide_height)

    center_x_norm = left_norm + (width_norm / 2.0)
    center_y_norm = top_norm + (height_norm / 2.0)

    return {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "left_norm": left_norm,
        "top_norm": top_norm,
        "width_norm": width_norm,
        "height_norm": height_norm,
        "center_x_norm": center_x_norm,
        "center_y_norm": center_y_norm,
        "area_norm": max(width_norm * height_norm, 0.0),
    }


def _get_placeholder_type(shape) -> Optional[Any]:
    try:
        if not shape.is_placeholder:
            return None
        return shape.placeholder_format.type
    except Exception:
        return None


def _is_title_placeholder_shape(shape) -> bool:
    placeholder_type = _get_placeholder_type(shape)
    return placeholder_type in TITLE_PLACEHOLDER_TYPES if placeholder_type is not None else False


def _shape_bounds(shape, slide_width: float, slide_height: float) -> Tuple[float, float, float, float]:
    position = _shape_position(shape, slide_width, slide_height)
    left = float(position.get("left_norm") or 0.0)
    top = float(position.get("top_norm") or 0.0)
    width = float(position.get("width_norm") or 0.0)
    height = float(position.get("height_norm") or 0.0)
    return left, top, left + width, top + height


def _is_hidden_slide(slide) -> bool:
    try:
      xml = slide._element.xml  # type: ignore[attr-defined]
      return 'show="0"' in xml or 'hidden="1"' in xml
    except Exception:
      return False


def _is_hidden_shape(shape) -> bool:
    try:
        xml = shape._element.xml  # type: ignore[attr-defined]
    except Exception:
        xml = ""
    lowered = xml.lower()
    return 'hidden="1"' in lowered or 'fhidden="1"' in lowered or 'show="0"' in lowered


def _is_off_slide_shape(shape, slide_width: float, slide_height: float) -> bool:
    left, top, right, bottom = _shape_bounds(shape, slide_width, slide_height)
    return right <= -0.02 or bottom <= -0.02 or left >= 1.02 or top >= 1.02


def _looks_like_placeholder_residue(text: str) -> bool:
    lowered = _normalize_pptx_text(text).lower()
    if not lowered:
        return False
    return any(marker in lowered for marker in PLACEHOLDER_RESIDUE_PATTERNS)


def _suppressed_shape(
    reason: str,
    risk_flag: Optional[str],
    preview: Optional[str],
    visibility_source: str,
) -> Dict[str, Any]:
    return {
        "type": "suppressed",
        "reason": reason,
        "risk_flag": risk_flag,
        "preview": preview,
        "visibility_source": visibility_source,
    }


def _extract_slide_title_fallback(slide, slide_width: float, slide_height: float) -> Optional[str]:
    """
    Recover title text for slides where python-pptx does not populate slide.shapes.title.

    Candidate policy:
    - only TITLE / CENTER_TITLE placeholders
    - choose smallest top_norm
    - tie-breaker: larger font size, then longer text (capped at 120 chars)
    """
    candidates: List[Tuple[float, float, int, str]] = []

    for shape in getattr(slide, "shapes", []):
        if not _is_title_placeholder_shape(shape):
            continue
        text = extract_text_from_shape(shape)
        text = _normalize_pptx_text(text)
        if not text:
            continue

        try:
            position = _shape_position(shape, slide_width, slide_height)
            top_norm = float(position.get("top_norm") or 1.0)
        except Exception:
            top_norm = 1.0

        try:
            text_meta = _extract_text_meta(shape)
            size = float(text_meta.get("font_size_max_pt") or text_meta.get("font_size_avg_pt") or 0.0)
        except Exception:
            size = 0.0

        length_score = min(len(text), 120)
        candidates.append((top_norm, -size, -length_score, text))

    if not candidates:
        return None

    candidates.sort(key=lambda item: (item[0], item[1], item[2]))
    return candidates[0][3]


def _extract_text_meta(shape) -> Dict[str, Any]:
    """Estimate headline-likeness via font size + weight."""
    try:
        if not shape.has_text_frame:
            return {}

        sizes: List[float] = []
        bold_runs = 0
        total_runs = 0

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                total_runs += 1
                font = run.font
                if font is not None:
                    if font.size is not None:
                        try:
                            sizes.append(float(font.size.pt))
                        except Exception:
                            pass
                    if font.bold:
                        bold_runs += 1

        max_size = max(sizes) if sizes else None
        avg_size = (sum(sizes) / len(sizes)) if sizes else None
        bold_ratio = (bold_runs / total_runs) if total_runs else 0.0

        return {
            "font_size_max_pt": max_size,
            "font_size_avg_pt": avg_size,
            "is_bold": bold_ratio >= 0.5 if total_runs else False,
            "bold_ratio": bold_ratio,
        }
    except Exception:
        return {}


def process_shape(
    shape,
    slide_num: int,
    slide_width: float,
    slide_height: float,
    *,
    include_images: bool = True,
) -> Optional[Dict[str, Any]]:
    """Process a single shape and extract relevant data."""
    try:
        # Skip footer shapes (date, slide number, footer text)
        if is_footer_shape(shape):
            return None

        if _is_hidden_shape(shape):
            preview = extract_text_from_shape(shape) if getattr(shape, "has_text_frame", False) else None
            return _suppressed_shape("hidden_shape", "hidden_text_leak", preview, "hidden_shape")

        if _is_off_slide_shape(shape, slide_width, slide_height):
            preview = extract_text_from_shape(shape) if getattr(shape, "has_text_frame", False) else None
            return _suppressed_shape("off_slide_shape", "off_slide_text", preview, "off_slide")

        # Image shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if not include_images:
                return None
            return extract_image(shape, slide_num, slide_width, slide_height)

        # Group shapes (may contain images)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            images = []
            texts = []
            for sub_shape in shape.shapes:
                if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE and include_images:
                    img = extract_image(sub_shape, slide_num, slide_width, slide_height)
                    if img:
                        images.append(img)
                if getattr(sub_shape, "has_text_frame", False):
                    text = extract_text_from_shape(sub_shape)
                    if text:
                        texts.append(text)
            if images or texts:
                group_data: Dict[str, Any] = {"type": "group"}
                if images:
                    group_data["images"] = images
                if texts:
                    group_data["text"] = "\n".join(texts)
                return group_data
            return None

        # Text shapes
        if shape.has_text_frame:
            text = extract_text_from_shape(shape)
            if text:
                if _looks_like_placeholder_residue(text):
                    return _suppressed_shape("placeholder_residue", "master_artifact", text, "placeholder_residue")
                position = _shape_position(shape, slide_width, slide_height)
                text_meta = _extract_text_meta(shape)
                placeholder_type = _get_placeholder_type(shape)
                is_title_shape = _is_title_placeholder_shape(shape)
                if not is_title_shape:
                    try:
                        slide_title_shape = shape.part.slide.shapes.title if hasattr(shape.part, "slide") else None
                        is_title_shape = bool(slide_title_shape is not None and shape == slide_title_shape)
                    except Exception:
                        is_title_shape = False
                return {
                    "type": "text",
                    "content": text,
                    "is_title": is_title_shape,
                    "placeholder_type": str(placeholder_type) if placeholder_type is not None else None,
                    "visibility_source": "visible_slide",
                    "classification_source": "ooxml",
                    **position,
                    **text_meta,
                }

        # Tables
        if shape.has_table:
            return extract_table(shape)

        # Charts
        if shape.has_chart:
            print(f"[PPTX Parser] Chart detected on slide {slide_num}!")
            return extract_chart(shape)

        # Other shapes (auto shapes, etc.) - include text if present
        if hasattr(shape, 'shape_type'):
            shape_data = {
                "type": "shape",
                "shape_type": str(shape.shape_type),
                "visibility_source": "visible_slide",
                "classification_source": "ooxml",
            }
            # Also extract text from shapes that have text frames (for classifier)
            if shape.has_text_frame:
                text = extract_text_from_shape(shape)
                if text:
                    shape_data["text"] = text
            return shape_data

    except Exception as e:
        print(f"[PPTX Parser] Error processing shape on slide {slide_num}: {e}")

    return None


ROW_TOLERANCE = 0.045


def _row_key(block: Dict[str, Any]) -> float:
    return float(block.get("top_norm") or 0.0)


def sort_text_blocks_by_layout(blocks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """
    Sort blocks into a closer reading order using row clustering.

    This helps preserve "card/grid" layouts where shape order is not semantic.
    """
    if not blocks:
        return []

    # Use normalized positions when available; otherwise keep original order.
    if not any("top_norm" in b for b in blocks):
        return list(blocks)

    rows: List[Dict[str, Any]] = []
    for block in sorted(blocks, key=_row_key):
        top = _row_key(block)
        assigned = False
        for row in rows:
            if abs(top - row["mean_top"]) <= ROW_TOLERANCE:
                row["blocks"].append(block)
                row["mean_top"] = sum(_row_key(b) for b in row["blocks"]) / max(len(row["blocks"]), 1)
                assigned = True
                break
        if not assigned:
            rows.append({"mean_top": top, "blocks": [block]})

    ordered: List[Dict[str, Any]] = []
    for row in sorted(rows, key=lambda r: r["mean_top"]):
        row_blocks = sorted(row["blocks"], key=lambda b: float(b.get("left_norm") or 0.0))
        ordered.extend(row_blocks)

    # Annotate reading order index without mutating the original references too much.
    for idx, block in enumerate(ordered):
        block["reading_order"] = idx

    return ordered


NUMBER_RE = re.compile(r"^\s*(\d{1,2})[.)]?\s*$")


def _is_number_block(block: Dict[str, Any]) -> Optional[str]:
    content = str(block.get("content") or "").strip()
    if not content or len(content) > 4:
        return None
    match = NUMBER_RE.match(content)
    if not match:
        return None
    return match.group(1)


def _split_text_lines(raw: Optional[str]) -> List[str]:
    if not raw:
        return []
    text = str(raw)
    # Normalize PPTX soft line breaks (vertical tabs) to real newlines.
    text = (
        text.replace("\r\n", "\n")
        .replace("\r", "\n")
        .replace("\v", "\n")
        .replace("\x0b", "\n")
        .replace("\x0c", "\n")
        .replace("\u2028", "\n")
        .replace("\u2029", "\n")
    )
    parts = [part.strip() for part in text.split("\n")]
    return [part for part in parts if part]


def _normalize_pptx_text(raw: Optional[str]) -> str:
    lines = _split_text_lines(raw)
    return "\n".join(lines).strip()


def _distance(a: Dict[str, Any], b: Dict[str, Any]) -> float:
    ax = float(a.get("center_x_norm") or 0.0)
    ay = float(a.get("center_y_norm") or 0.0)
    bx = float(b.get("center_x_norm") or 0.0)
    by = float(b.get("center_y_norm") or 0.0)
    dx = ax - bx
    dy = ay - by
    return (dx * dx + dy * dy) ** 0.5


def _region_bounds(num_block: Dict[str, Any], other_numbers: List[Dict[str, Any]]) -> Tuple[float, float, float, float]:
    """Create a loose bounding box around a numbered item."""
    left = float(num_block.get("left_norm") or 0.0)
    top = float(num_block.get("top_norm") or 0.0)
    width = float(num_block.get("width_norm") or 0.0)
    height = float(num_block.get("height_norm") or 0.0)
    cx = float(num_block.get("center_x_norm") or (left + width / 2.0))
    cy = float(num_block.get("center_y_norm") or (top + height / 2.0))

    min_x = max(0.0, left - 0.03)
    max_x = min(1.0, left + max(0.45, width * 4.0 + 0.25))
    min_y = max(0.0, top - 0.06)
    max_y = min(1.0, top + max(0.38, height * 6.0 + 0.25))

    # Clamp the region vertically between nearby number blocks in the same band.
    below = [
        nb for nb in other_numbers
        if float(nb.get("center_y_norm") or 0.0) > cy
        and abs(float(nb.get("center_x_norm") or 0.0) - cx) < 0.28
    ]
    if below:
        next_nb = min(below, key=lambda nb: float(nb.get("center_y_norm") or 1.0))
        next_cy = float(next_nb.get("center_y_norm") or max_y)
        max_y = min(max_y, cy + ((next_cy - cy) * 0.6))

    return min_x, max_x, min_y, max_y


def _in_bounds(block: Dict[str, Any], bounds: Tuple[float, float, float, float]) -> bool:
    min_x, max_x, min_y, max_y = bounds
    cx = float(block.get("center_x_norm") or 0.0)
    cy = float(block.get("center_y_norm") or 0.0)
    return min_x <= cx <= max_x and min_y <= cy <= max_y


def detect_numbered_groups(blocks: List[Dict[str, Any]]) -> Tuple[List[Dict[str, Any]], Set[int]]:
    """
    Detect layouts like "1 + heading + text" repeated across a slide.

    Returns structured items and the set of source indices that were grouped.
    """
    number_blocks = []
    for block in blocks:
        num = _is_number_block(block)
        if num:
            number_blocks.append((block, num))

    # Require at least 2 numbered anchors to reduce false positives.
    if len(number_blocks) < 2:
        return [], set()

    structured_items: List[Dict[str, Any]] = []
    used_indices: Set[int] = set()
    number_only = [nb for nb, _ in number_blocks]

    for num_block, num_value in number_blocks:
        bounds = _region_bounds(num_block, number_only)
        candidates = [
            b for b in blocks
            if b is not num_block
            and not _is_number_block(b)
            and _in_bounds(b, bounds)
            and isinstance(b.get("index"), int)
            and int(b.get("index")) not in used_indices
        ]

        if not candidates:
            continue

        # Prefer nearby, right-of-number, and typographically larger blocks as headings.
        def heading_score(block: Dict[str, Any]) -> float:
            dist = _distance(num_block, block)
            cx = float(block.get("center_x_norm") or 0.0)
            nx = float(num_block.get("center_x_norm") or 0.0)
            right_bonus = -0.03 if cx >= nx - 0.01 else 0.02
            font = float(block.get("font_size_max_pt") or 0.0)
            font_bonus = min(font, 48.0) / 400.0
            return dist + right_bonus - font_bonus

        heading = min(candidates, key=heading_score)
        heading_idx = int(heading.get("index"))

        heading_lines = _split_text_lines(heading.get("content"))
        heading_text = heading_lines[0] if heading_lines else str(heading.get("content") or "").strip()
        heading_body_lines = heading_lines[1:]

        remaining = [c for c in candidates if c is not heading]
        remaining.sort(key=lambda b: float(b.get("top_norm") or 0.0))

        body_blocks = []
        heading_y = float(heading.get("center_y_norm") or 0.0)
        for block in remaining:
            if float(block.get("center_y_norm") or 0.0) + 0.01 >= heading_y:
                body_blocks.append(block)
            if len(body_blocks) >= 3:
                break

        body_texts: List[str] = []
        for b in body_blocks:
            if b.get("content"):
                body_texts.extend(_split_text_lines(b.get("content")))
        if heading_body_lines:
            body_texts = heading_body_lines + body_texts

        source_indices = {int(num_block.get("index"))} if isinstance(num_block.get("index"), int) else set()
        source_indices.add(heading_idx)
        for b in body_blocks:
            if isinstance(b.get("index"), int):
                source_indices.add(int(b.get("index")))

        used_indices.update(source_indices)

        structured_items.append(
            {
                "number": num_value,
                "heading": heading_text,
                "body": body_texts,
                "source_indices": sorted(source_indices),
            }
        )

    # Sort items by their visual order (top, then left).
    structured_items.sort(
        key=lambda item: (
            min(
                (
                    float(next((b.get("top_norm") for b in blocks if b.get("index") == idx), 0.0))
                    for idx in item.get("source_indices", [])
                ),
                default=0.0,
            ),
            min(
                (
                    float(next((b.get("left_norm") for b in blocks if b.get("index") == idx), 0.0))
                    for idx in item.get("source_indices", [])
                ),
                default=0.0,
            ),
        )
    )

    return structured_items, used_indices


def build_structured_text(structured_items: List[Dict[str, Any]]) -> str:
    """Render structured items into a compact, LLM-friendly text block."""
    lines: List[str] = []
    for item in structured_items:
        number = str(item.get("number") or "").strip()
        heading = str(item.get("heading") or "").strip()
        body = " ".join(str(t).strip() for t in item.get("body") or [] if str(t).strip())

        prefix = f"{number}. " if number else ""
        if heading and body:
            lines.append(f"{prefix}{heading}: {body}".strip())
        elif heading:
            lines.append(f"{prefix}{heading}".strip())
        elif body:
            lines.append(f"{prefix}{body}".strip())

    return "\n".join(line for line in lines if line)


def extract_image(
    shape,
    slide_num: int,
    slide_width: float,
    slide_height: float,
) -> Optional[Dict[str, Any]]:
    """Extract image data from a picture shape."""
    try:
        image = shape.image
        image_bytes = image.blob

        # Get image format
        content_type = image.content_type
        ext = content_type.split("/")[-1] if content_type else "png"

        # Get dimensions
        width = shape.width.inches if shape.width else 0
        height = shape.height.inches if shape.height else 0
        position = _shape_position(shape, slide_width, slide_height)

        # Check if image has existing alt text
        alt_text = None
        if hasattr(shape, '_element'):
            # Try to get alt text from shape properties
            nvPicPr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr')
            if nvPicPr is not None:
                cNvPr = nvPicPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
                if cNvPr is not None:
                    alt_text = cNvPr.get('descr')

        # Create base64 data URI for HTML embedding
        b64_data = base64.b64encode(image_bytes).decode('utf-8')
        data_uri = f"data:{content_type};base64,{b64_data}"

        return {
            "type": "image",
            "image_bytes": image_bytes,
            "data_uri": data_uri,
            "content_type": content_type,
            "extension": ext,
            "width_inches": width,
            "height_inches": height,
            "existing_alt_text": alt_text,
            "alt_text": alt_text,  # Will be updated by VLM if needed
            "slide_number": slide_num,
            "visibility_source": "visible_slide",
            "classification_source": "ooxml",
            **position,
        }

    except Exception as e:
        print(f"[PPTX Parser] Error extracting image on slide {slide_num}: {e}")
        return None


def extract_text_from_shape(shape) -> Optional[str]:
    """Extract text content from a text frame."""
    try:
        if not shape.has_text_frame:
            return None

        paragraphs = []
        for para in shape.text_frame.paragraphs:
            text = _normalize_pptx_text(para.text)
            if text:
                paragraphs.append(text)

        return "\n".join(paragraphs) if paragraphs else None

    except Exception as e:
        print(f"[PPTX Parser] Error extracting text: {e}")
        return None


def extract_table(shape) -> Dict[str, Any]:
    """Extract table data from a table shape, including merged-cell info and column widths."""
    try:
        table = shape.table
        num_rows = len(table.rows)
        num_cols = len(table.columns)

        # --- column widths as relative fractions (0-1) -----------------------
        col_widths: list = []
        try:
            raw_widths = [col.width for col in table.columns]
            total = sum(raw_widths) or 1
            col_widths = [round(w / total, 4) for w in raw_widths]
        except Exception:
            col_widths = []

        # --- build a grid to track which cells are "covered" by a merge ------
        covered: set = set()  # (row_idx, col_idx) pairs swallowed by a span

        rows: list = []
        for row_idx, row in enumerate(table.rows):
            cells: list = []
            for col_idx, cell in enumerate(row.cells):
                # Skip continuation cells that are covered by a prior merge
                if (row_idx, col_idx) in covered:
                    continue

                cell_text = cell.text.strip() if cell.text else ""
                colspan = 1
                rowspan = 1

                # Detect merge origin
                try:
                    if cell.is_merge_origin:
                        colspan = cell.span_width
                        rowspan = cell.span_height
                        # Mark all spanned cells (except the origin) as covered
                        for dr in range(rowspan):
                            for dc in range(colspan):
                                if dr == 0 and dc == 0:
                                    continue
                                covered.add((row_idx + dr, col_idx + dc))
                except Exception:
                    # python-pptx < 0.6.22 may not expose these attributes;
                    # fall back to no-merge behaviour.
                    pass

                cells.append({
                    "text": cell_text,
                    "colspan": colspan,
                    "rowspan": rowspan,
                })
            rows.append(cells)

        return {
            "type": "table",
            "classification_source": "ooxml",
            "rows": rows,
            "col_widths": col_widths,
            "row_count": num_rows,
            "col_count": num_cols,
        }

    except Exception as e:
        print(f"[PPTX Parser] Error extracting table: {e}")
        return {
            "type": "table",
            "classification_source": "ooxml",
            "rows": [],
            "col_widths": [],
            "row_count": 0,
            "col_count": 0,
        }


def extract_chart(shape) -> Dict[str, Any]:
    """Extract chart data for generating accessible description."""
    try:
        chart = shape.chart
        chart_type = str(chart.chart_type) if chart else "unknown"

        # Get title
        title = None
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip()

        # Extract series data for description
        series_data = []
        try:
            for series in chart.series:
                series_info = {
                    "name": series.name if hasattr(series, 'name') else None,
                    "values": [],
                }

                # Try to extract values
                if hasattr(series, 'values') and series.values:
                    series_info["values"] = list(series.values)[:10]  # Limit to 10 values

                series_data.append(series_info)
        except Exception:
            pass

        # Extract categories (x-axis labels)
        categories = []
        try:
            if hasattr(chart, 'plots') and chart.plots:
                plot = chart.plots[0]
                if hasattr(plot, 'categories') and plot.categories:
                    categories = list(plot.categories)[:10]  # Limit to 10
        except Exception:
            pass

        # Map chart type to German description
        chart_type_german = get_chart_type_german(chart_type)

        return {
            "type": "chart",
            "classification_source": "ooxml",
            "chart_type": chart_type,
            "chart_type_german": chart_type_german,
            "title": title,
            "series_data": series_data,
            "categories": categories,
            "summary": None,  # Will be filled by generate_chart_summary
        }

    except Exception as e:
        print(f"[PPTX Parser] Error extracting chart: {e}")
        return {
            "type": "chart",
            "classification_source": "ooxml",
            "chart_type": "unknown",
            "chart_type_german": "Diagramm",
            "title": None,
            "series_data": [],
            "categories": [],
            "summary": None,
        }


def get_chart_type_german(chart_type: str) -> str:
    """Map PowerPoint chart type to German name."""
    type_map = {
        "BAR_CLUSTERED": "Balkendiagramm",
        "BAR_STACKED": "Gestapeltes Balkendiagramm",
        "COLUMN_CLUSTERED": "Säulendiagramm",
        "COLUMN_STACKED": "Gestapeltes Säulendiagramm",
        "LINE": "Liniendiagramm",
        "LINE_MARKERS": "Liniendiagramm mit Markern",
        "PIE": "Kreisdiagramm",
        "DOUGHNUT": "Ringdiagramm",
        "AREA": "Flächendiagramm",
        "AREA_STACKED": "Gestapeltes Flächendiagramm",
        "SCATTER": "Punktdiagramm",
        "BUBBLE": "Blasendiagramm",
        "RADAR": "Netzdiagramm",
    }
    return type_map.get(chart_type, "Diagramm")
