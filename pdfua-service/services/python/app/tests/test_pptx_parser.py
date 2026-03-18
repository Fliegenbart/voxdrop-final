from types import SimpleNamespace
import sys

sys.modules.setdefault("pdf2image", SimpleNamespace(convert_from_path=lambda *args, **kwargs: []))
sys.modules.setdefault("pypdf", SimpleNamespace(PdfReader=object))

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from app.parsers.pptx_parser import _extract_slide_title_fallback, process_shape


class _FakeParagraph:
    def __init__(self, text: str, size_pt: float = 24, bold: bool = True):
        self.text = text
        self.runs = [
            SimpleNamespace(
                font=SimpleNamespace(
                    size=SimpleNamespace(pt=size_pt) if size_pt else None,
                    bold=bold,
                )
            )
        ]


class _FakeTextFrame:
    def __init__(self, paragraphs):
        self.paragraphs = paragraphs


class _FakePlaceholderFormat:
    def __init__(self, placeholder_type):
        self.type = placeholder_type


class _FakeSlideRef:
    def __init__(self, title_shape=None):
        self.shapes = SimpleNamespace(title=title_shape)


class _FakePart:
    def __init__(self, title_shape=None):
        self.slide = _FakeSlideRef(title_shape=title_shape)


class _FakeShape:
    def __init__(
        self,
        text: str,
        *,
        placeholder_type=PP_PLACEHOLDER.TITLE,
        top: int = 1000,
        left: int = 1000,
        width: int = 4000,
        height: int = 900,
        size_pt: float = 32,
        title_shape=None,
        xml: str = "<p:sp/>",
    ):
        self.has_text_frame = True
        self.text_frame = _FakeTextFrame([_FakeParagraph(text, size_pt=size_pt, bold=True)])
        self.is_placeholder = True
        self.placeholder_format = _FakePlaceholderFormat(placeholder_type)
        self.top = top
        self.left = left
        self.width = width
        self.height = height
        self.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        self.has_table = False
        self.has_chart = False
        self.part = _FakePart(title_shape=title_shape)
        self._element = SimpleNamespace(xml=xml)


def test_parser_recovers_title_when_slide_shapes_title_is_none_but_title_placeholder_exists():
    slide = SimpleNamespace(
        shapes=[
            _FakeShape("Unsere Ziele bis Ende 2029", top=900, size_pt=36),
            _FakeShape("AGILE", top=2200, size_pt=20),
        ]
    )

    title = _extract_slide_title_fallback(slide, slide_width=10000.0, slide_height=10000.0)
    assert title == "Unsere Ziele bis Ende 2029"


def test_process_shape_marks_title_placeholder_as_is_title():
    shape = _FakeShape("Unsere Ziele bis Ende 2029", placeholder_type=PP_PLACEHOLDER.CENTER_TITLE, title_shape=None)

    parsed = process_shape(
        shape,
        slide_num=7,
        slide_width=10000.0,
        slide_height=10000.0,
        include_images=True,
    )

    assert parsed is not None
    assert parsed["type"] == "text"
    assert parsed["is_title"] is True
    assert parsed["content"] == "Unsere Ziele bis Ende 2029"


def test_process_shape_suppresses_hidden_text_shapes():
    shape = _FakeShape("Nur intern", xml='<p:sp hidden="1"/>')

    parsed = process_shape(
        shape,
        slide_num=3,
        slide_width=10000.0,
        slide_height=10000.0,
        include_images=True,
    )

    assert parsed is not None
    assert parsed["type"] == "suppressed"
    assert parsed["risk_flag"] == "hidden_text_leak"


def test_process_shape_suppresses_placeholder_residue():
    shape = _FakeShape("Klicken Sie, um Text hinzuzufügen", placeholder_type=PP_PLACEHOLDER.BODY)

    parsed = process_shape(
        shape,
        slide_num=4,
        slide_width=10000.0,
        slide_height=10000.0,
        include_images=True,
    )

    assert parsed is not None
    assert parsed["type"] == "suppressed"
    assert parsed["risk_flag"] == "master_artifact"
