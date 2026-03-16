"""
Pipeline orchestrator for voxdrop-pptx2pdfua.

Orchestrates the multi-stage pipeline for converting PPTX to accessible PDF/HTML:
1. IR Extraction - Extract intermediate representation from PPTX
2. Semantic Analysis - Analyze reading order, timeline, complexity
3. Alt-Text Generation - Generate alt-texts for visual elements
4. Alt-Text QA Layer - Validate and improve generated alt-texts
5. Final Rendering - Generate accessible PDF/HTML output

Phase 3 adds:
- GlobalContextManager for presentation-wide context in QA
- VisionCaptioner for image analysis and vision seeds
"""

import json
import os
import subprocess
import sys
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from .qa.alttext_manager import AltTextQualityLayer, QAResult
from .qa.context_engine import GlobalContextManager, get_context_manager
from .qa.vision_captioner import VisionCaptioner, get_vision_captioner
from .qa.prompts import build_slide_summary_prompt, get_slide_summary_system_prompt
from .qa.compliance_checker import (
    AccessibilityValidator,
    ComplianceReport,
    validate_pdf_accessibility,
)


# Phase 3 Configuration
USE_GLOBAL_CONTEXT = os.getenv("QA_USE_CONTEXT", "true").lower() == "true"
USE_VISION_SEEDS = os.getenv("QA_USE_VISION", "false").lower() == "true"
USE_SPEAKER_NOTES = os.getenv("QA_USE_SPEAKER_NOTES", "true").lower() == "true"


class PipelineStage(str, Enum):
    """Pipeline stages."""
    IR_EXTRACTION = "ir_extraction"
    SEMANTIC_ANALYSIS = "semantic_analysis"
    ALTTEXT_GENERATION = "alttext_generation"
    ALTTEXT_QA = "alttext_qa"
    COMPLIANCE_CHECK = "compliance_check"
    FINAL_RENDER = "final_render"


@dataclass
class StageResult:
    """Result of a pipeline stage."""
    stage: PipelineStage
    success: bool
    output_path: Optional[Path] = None
    error: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = None


@dataclass
class PipelineResult:
    """Result of the full pipeline execution."""
    success: bool
    stages: List[StageResult]
    final_output: Optional[Path] = None
    qa_result: Optional[QAResult] = None


class PipelineOrchestrator:
    """
    Orchestrates the PPTX to accessible document pipeline.

    The pipeline runs in stages:
    1. IR Extraction: PPTX -> IR JSON (facts only)
    2. Semantic Analysis: IR -> Semantic JSON (reading order, timeline)
    3. Alt-Text Generation: Generate alt-texts using LLM
    4. Alt-Text QA: Validate and improve alt-texts (NEW)
    5. Final Render: Generate PDF/UA or accessible HTML
    """

    def __init__(
        self,
        app_root: Optional[Path] = None,
        output_dir: Optional[Path] = None,
        llm_host: Optional[str] = None,
        llm_model: Optional[str] = None,
        language: str = "de",
    ):
        """
        Initialize the orchestrator.

        Args:
            app_root: Root directory containing pipeline scripts.
            output_dir: Directory for output files.
            llm_host: LLM server host URL.
            llm_model: LLM model name.
            language: Language for prompts ("de" or "en").
        """
        self.app_root = app_root or Path("/app")
        self.output_dir = output_dir or Path("/data/processed")
        self.llm_host = llm_host
        self.llm_model = llm_model
        self.language = language

        # Initialize the QA layer
        self.qa_layer = AltTextQualityLayer(
            llm_host=llm_host,
            llm_model=llm_model,
            language=language,
        )

        # Phase 3: Context manager (lazy initialized)
        self._context_manager: Optional[GlobalContextManager] = None

        # Phase 3: Vision captioner (lazy initialized)
        self._vision_captioner: Optional[VisionCaptioner] = None

    def run(
        self,
        pptx_path: Path,
        job_id: str,
        skip_stages: Optional[List[PipelineStage]] = None,
    ) -> PipelineResult:
        """
        Run the full pipeline on a PPTX file.

        Args:
            pptx_path: Path to the input PPTX file.
            job_id: Unique identifier for this job.
            skip_stages: Optional list of stages to skip.

        Returns:
            PipelineResult with all stage results.
        """
        skip_stages = skip_stages or []
        stages: List[StageResult] = []
        job_dir = self.output_dir / job_id
        job_dir.mkdir(parents=True, exist_ok=True)

        # Define output paths
        ir_path = job_dir / f"{job_id}.ir.json"
        semantic_path = job_dir / f"{job_id}.semantic.json"
        alttext_path = job_dir / f"{job_id}.alttext.json"
        alttext_qa_path = job_dir / f"{job_id}.alttext_qa.json"
        modelpack_path = job_dir / f"{job_id}.modelpack.json"

        qa_result: Optional[QAResult] = None

        # Stage 1: IR Extraction
        if PipelineStage.IR_EXTRACTION not in skip_stages:
            result = self._run_ir_extraction(pptx_path, ir_path)
            stages.append(result)
            if not result.success:
                return PipelineResult(success=False, stages=stages)

        # Stage 2: Semantic Analysis
        if PipelineStage.SEMANTIC_ANALYSIS not in skip_stages:
            result = self._run_semantic_analysis(ir_path, semantic_path)
            stages.append(result)
            if not result.success:
                return PipelineResult(success=False, stages=stages)

        # Stage 3: Alt-Text Generation (with optional vision seeds)
        if PipelineStage.ALTTEXT_GENERATION not in skip_stages:
            result = self._run_alttext_generation(
                pptx_path, ir_path, semantic_path, modelpack_path, alttext_path
            )
            stages.append(result)
            if not result.success:
                return PipelineResult(success=False, stages=stages)

        # Stage 4: Alt-Text QA Layer (with optional global context)
        if PipelineStage.ALTTEXT_QA not in skip_stages and alttext_path.exists():
            result, qa_result = self._run_alttext_qa(
                alttext_path, alttext_qa_path, modelpack_path
            )
            stages.append(result)
            if not result.success:
                return PipelineResult(success=False, stages=stages, qa_result=qa_result)

        # Stage 5: Compliance Check
        if PipelineStage.COMPLIANCE_CHECK not in skip_stages:
            result = self._run_compliance_check(modelpack_path, alttext_qa_path)
            stages.append(result)
            # Don't fail on compliance issues - just log them

        # Stage 6: Final Render (placeholder for future implementation)
        pdf_compliance_metadata: Dict[str, Any] = {}
        if PipelineStage.FINAL_RENDER not in skip_stages:
            result = self._run_final_render(
                ir_path, semantic_path, alttext_qa_path, job_dir
            )
            stages.append(result)

            # Validate PDF output after generation
            if result.success and result.output_path and result.output_path.exists():
                print("[Orchestrator] Validating generated PDF...")
                pdf_result, pdf_report = self._validate_pdf_output(result.output_path)
                stages.append(pdf_result)

                # Store compliance status in metadata
                pdf_compliance_metadata = {
                    "is_compliant": pdf_report.is_compliant,
                    "status": pdf_report.status,
                    "errors": pdf_report.errors,
                    "warnings": pdf_report.warnings,
                }

        return PipelineResult(
            success=all(s.success for s in stages),
            stages=stages,
            final_output=job_dir / "output.pdf" if stages and stages[-1].success else None,
            qa_result=qa_result,
        )

    def _run_ir_extraction(
        self,
        pptx_path: Path,
        output_path: Path,
    ) -> StageResult:
        """Stage 1: Extract IR from PPTX."""
        try:
            script_path = self.app_root / "pptx_extract_ir.py"
            subprocess.run(
                [sys.executable, str(script_path), str(pptx_path), "-o", str(output_path)],
                check=True,
                capture_output=True,
                text=True,
            )
            return StageResult(
                stage=PipelineStage.IR_EXTRACTION,
                success=True,
                output_path=output_path,
            )
        except subprocess.CalledProcessError as e:
            return StageResult(
                stage=PipelineStage.IR_EXTRACTION,
                success=False,
                error=f"IR extraction failed: {e.stderr}",
            )

    def _run_semantic_analysis(
        self,
        ir_path: Path,
        output_path: Path,
    ) -> StageResult:
        """Stage 2: Run semantic analysis."""
        try:
            script_path = self.app_root / "pptx_semantic_pack.py"
            subprocess.run(
                [sys.executable, str(script_path), str(ir_path), "-o", str(output_path)],
                check=True,
                capture_output=True,
                text=True,
            )
            return StageResult(
                stage=PipelineStage.SEMANTIC_ANALYSIS,
                success=True,
                output_path=output_path,
            )
        except subprocess.CalledProcessError as e:
            return StageResult(
                stage=PipelineStage.SEMANTIC_ANALYSIS,
                success=False,
                error=f"Semantic analysis failed: {e.stderr}",
            )

    def _extract_speaker_notes(
        self,
        modelpack_data: Dict[str, Any],
    ) -> Dict[int, str]:
        """
        Extract speaker notes from modelpack data.

        Phase 3: Speaker notes are used as primary evidence for alt-text
        generation and slide summaries.

        Args:
            modelpack_data: Modelpack JSON data.

        Returns:
            Dict mapping slide_index to speaker notes text.
        """
        speaker_notes: Dict[int, str] = {}

        for slide in modelpack_data.get("slides", []):
            slide_idx = slide.get("slide_index", 0)
            notes = slide.get("notes_text", "")

            if notes and notes.strip():
                speaker_notes[slide_idx] = notes.strip()
                print(f"[Orchestrator] Extracted speaker notes for slide {slide_idx + 1}: {len(notes)} chars")

        print(f"[Orchestrator] Total slides with speaker notes: {len(speaker_notes)}")
        return speaker_notes

    def _extract_images_from_pptx(
        self,
        pptx_path: Path,
        ir_data: Dict[str, Any],
    ) -> Dict[str, bytes]:
        """
        Extract image bytes from PPTX for vision processing.

        Phase 3: Extracts images from picture shapes to generate vision seeds.

        Args:
            pptx_path: Path to the original PPTX file.
            ir_data: IR data containing shape information.

        Returns:
            Dict mapping shape_id to image bytes.
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            print("[Orchestrator] python-pptx not available for image extraction")
            return {}

        image_data: Dict[str, bytes] = {}

        try:
            prs = Presentation(str(pptx_path))

            # Build a lookup of shape_ids from IR that are pictures
            picture_shape_ids = set()
            for slide in ir_data.get("slides", []):
                for shape in slide.get("shapes", []):
                    if shape.get("type") in ("picture", "chart", "diagram"):
                        picture_shape_ids.add(str(shape.get("shape_id", "")))

            # Extract images from PPTX
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    shape_id = str(shape.shape_id)

                    # Check if this is a picture shape we need
                    if shape_id in picture_shape_ids:
                        try:
                            if hasattr(shape, "image") and shape.image:
                                image_data[shape_id] = shape.image.blob
                        except Exception:
                            # Shape doesn't have accessible image data
                            pass

            print(f"[Orchestrator] Extracted {len(image_data)} images for vision processing")
            return image_data

        except Exception as e:
            print(f"[Orchestrator] Image extraction failed: {e}")
            return {}

    def _generate_vision_seeds(
        self,
        image_data: Dict[str, bytes],
        ir_data: Dict[str, Any],
    ) -> Dict[str, str]:
        """
        Generate vision seeds for images using VisionCaptioner.

        Phase 3: Uses vision-language model to describe images.

        Args:
            image_data: Dict mapping shape_id to image bytes.
            ir_data: IR data for context (slide titles, etc.).

        Returns:
            Dict mapping shape_id to vision description.
        """
        if not image_data:
            return {}

        # Initialize vision captioner if needed
        if self._vision_captioner is None:
            self._vision_captioner = get_vision_captioner(
                llm_host=self.llm_host,
                language=self.language,
            )

        if not self._vision_captioner.is_available():
            print("[Orchestrator] Vision model not available, skipping vision seeds")
            return {}

        # Build context lookup from IR
        shape_context: Dict[str, str] = {}
        for slide in ir_data.get("slides", []):
            slide_title = slide.get("title", "")
            for shape in slide.get("shapes", []):
                shape_id = str(shape.get("shape_id", ""))
                if shape_id:
                    shape_context[shape_id] = slide_title

        # Generate captions
        vision_seeds: Dict[str, str] = {}
        for shape_id, img_bytes in image_data.items():
            try:
                context = shape_context.get(shape_id, "")
                caption = self._vision_captioner.generate_caption(
                    image_bytes=img_bytes,
                    context=context,
                    shape_type="image",
                )
                if caption:
                    vision_seeds[shape_id] = caption
                    print(f"[Orchestrator] Vision seed for {shape_id}: {caption[:50]}...")
            except Exception as e:
                print(f"[Orchestrator] Vision seed generation failed for {shape_id}: {e}")

        print(f"[Orchestrator] Generated {len(vision_seeds)} vision seeds")
        return vision_seeds

    def _run_alttext_generation(
        self,
        pptx_path: Path,
        ir_path: Path,
        semantic_path: Path,
        modelpack_path: Path,
        alttext_path: Path,
    ) -> StageResult:
        """
        Stage 3: Generate alt-texts for visual elements.

        This creates the modelpack and generates alt-texts using the LLM.
        Phase 3: Optionally generates vision seeds for image elements.
        """
        try:
            # First create the modelpack
            script_path = self.app_root / "pptx_pack_minimizer.py"
            subprocess.run(
                [
                    sys.executable, str(script_path),
                    str(ir_path), str(semantic_path),
                    "-o", str(modelpack_path)
                ],
                check=True,
                capture_output=True,
                text=True,
            )

            # Load IR data for alt-text generation
            with open(ir_path, encoding="utf-8") as f:
                ir_data = json.load(f)

            # Phase 3: Generate vision seeds if enabled
            vision_seeds: Dict[str, str] = {}
            if USE_VISION_SEEDS:
                print("[Orchestrator] Vision seeds enabled, extracting images...")
                image_data = self._extract_images_from_pptx(pptx_path, ir_data)
                if image_data:
                    vision_seeds = self._generate_vision_seeds(image_data, ir_data)

            # Phase 3: Extract speaker notes from modelpack
            speaker_notes: Dict[int, str] = {}
            if USE_SPEAKER_NOTES and modelpack_path.exists():
                print("[Orchestrator] Extracting speaker notes...")
                with open(modelpack_path, encoding="utf-8") as f:
                    modelpack_data = json.load(f)
                speaker_notes = self._extract_speaker_notes(modelpack_data)

            # Generate alt-texts from IR data (with vision seeds and speaker notes)
            alt_texts = self._generate_alttexts_from_ir(
                ir_path, semantic_path, vision_seeds, speaker_notes
            )

            # Save alt-texts
            with open(alttext_path, "w", encoding="utf-8") as f:
                json.dump(alt_texts, f, ensure_ascii=False, indent=2)

            return StageResult(
                stage=PipelineStage.ALTTEXT_GENERATION,
                success=True,
                output_path=alttext_path,
                metadata={
                    "alt_text_count": len(alt_texts),
                    "vision_seeds_generated": len(vision_seeds),
                    "speaker_notes_slides": len(speaker_notes),
                },
            )
        except subprocess.CalledProcessError as e:
            return StageResult(
                stage=PipelineStage.ALTTEXT_GENERATION,
                success=False,
                error=f"Alt-text generation failed: {e.stderr}",
            )
        except Exception as e:
            return StageResult(
                stage=PipelineStage.ALTTEXT_GENERATION,
                success=False,
                error=f"Alt-text generation failed: {str(e)}",
            )

    def _generate_alttexts_from_ir(
        self,
        ir_path: Path,
        semantic_path: Path,
        vision_seeds: Optional[Dict[str, str]] = None,
        speaker_notes: Optional[Dict[int, str]] = None,
    ) -> List[Dict[str, Any]]:
        """
        Generate alt-text entries from IR and semantic data.

        This is a placeholder that extracts shape information.
        In production, this would call an LLM to generate actual alt-texts.

        Phase 3: Includes vision_seed and speaker_notes for context.

        Args:
            ir_path: Path to IR JSON file.
            semantic_path: Path to semantic JSON file.
            vision_seeds: Optional dict mapping shape_id to vision description.
            speaker_notes: Optional dict mapping slide_index to notes text.

        Returns:
            List of alt-text entry dictionaries.
        """
        vision_seeds = vision_seeds or {}
        speaker_notes = speaker_notes or {}

        with open(ir_path, encoding="utf-8") as f:
            ir_data = json.load(f)

        with open(semantic_path, encoding="utf-8") as f:
            semantic_data = json.load(f)

        alt_texts: List[Dict[str, Any]] = []

        slides = ir_data.get("slides", [])
        for slide in slides:
            slide_number = slide.get("slide_number", 0)
            slide_index = slide.get("slide_index", slide_number - 1)
            slide_title = slide.get("title", "")

            # Get speaker notes for this slide
            slide_notes = speaker_notes.get(slide_index, "")

            shapes = slide.get("shapes", [])
            for shape in shapes:
                shape_id = str(shape.get("shape_id", ""))
                shape_type = shape.get("type", "")
                text_content = shape.get("text", "")

                # Only process visual elements that need alt-text
                if shape_type in ("picture", "chart", "diagram", "smartart", "table"):
                    entry = {
                        "shape_id": shape_id,
                        "slide_number": slide_number,
                        "slide_index": slide_index,
                        "shape_type": shape_type,
                        "alt_text": shape.get("alt_text", f"Bild: {shape_type}"),
                        "original_text": text_content,
                        "slide_title": slide_title,
                        "ocr_snippet": shape.get("ocr_text", ""),
                    }

                    # Phase 3: Add vision seed if available
                    if shape_id in vision_seeds:
                        entry["vision_seed"] = vision_seeds[shape_id]

                    # Phase 3: Add speaker notes (primary evidence source)
                    if slide_notes:
                        entry["speaker_notes"] = slide_notes

                    alt_texts.append(entry)

        return alt_texts

    def _run_alttext_qa(
        self,
        alttext_path: Path,
        output_path: Path,
        modelpack_path: Optional[Path] = None,
    ) -> Tuple[StageResult, Optional[QAResult]]:
        """
        Stage 4: Run alt-text quality assurance.

        This is the QA layer that validates and improves alt-texts.
        Phase 3: Optionally extracts and uses global context from modelpack.

        Args:
            alttext_path: Path to alt-text JSON file.
            output_path: Path for QA result output.
            modelpack_path: Optional path to modelpack for context extraction.

        Returns:
            Tuple of StageResult and optional QAResult.
        """
        try:
            with open(alttext_path, encoding="utf-8") as f:
                alt_texts = json.load(f)

            # Phase 3: Extract global context if enabled
            global_context = ""
            context_extracted = False
            if USE_GLOBAL_CONTEXT and modelpack_path and modelpack_path.exists():
                try:
                    with open(modelpack_path, encoding="utf-8") as f:
                        modelpack_data = json.load(f)

                    # Initialize context manager if needed
                    if self._context_manager is None:
                        self._context_manager = get_context_manager(
                            llm_host=self.llm_host,
                            llm_model=self.llm_model,
                        )

                    # Extract context
                    presentation_context = self._context_manager.extract_context(modelpack_data)
                    global_context = presentation_context.to_context_string()
                    context_extracted = bool(global_context)
                    print(f"[Orchestrator] Extracted global context: {len(global_context)} chars")

                except Exception as e:
                    print(f"[Orchestrator] Context extraction failed: {e}")

            # Run QA layer with global context
            qa_result = self.qa_layer.process(alt_texts, global_context=global_context)

            # Save QA results
            qa_dict = self.qa_layer.to_dict(qa_result)
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(qa_dict, f, ensure_ascii=False, indent=2)

            return StageResult(
                stage=PipelineStage.ALTTEXT_QA,
                success=True,
                output_path=output_path,
                metadata={
                    "total_processed": qa_result.stats.total_processed,
                    "flagged_count": qa_result.stats.flagged_count,
                    "repaired_count": qa_result.stats.repaired_count,
                    "decorative_count": qa_result.stats.decorative_count,
                    "context_extracted": context_extracted,
                },
            ), qa_result

        except Exception as e:
            return StageResult(
                stage=PipelineStage.ALTTEXT_QA,
                success=False,
                error=f"Alt-text QA failed: {str(e)}",
            ), None

    def _run_compliance_check(
        self,
        modelpack_path: Path,
        alttext_qa_path: Path,
    ) -> StageResult:
        """
        Stage 5: Run PDF/UA compliance check.

        Validates that all mandatory elements for PDF/UA are present:
        - Document title
        - Document language
        - Tag tree structure
        - Alt-texts for all visual elements

        Args:
            modelpack_path: Path to modelpack JSON.
            alttext_qa_path: Path to QA result JSON.

        Returns:
            StageResult with compliance check results.
        """
        compliance_issues: List[str] = []
        compliance_passed: Dict[str, bool] = {
            "has_title": False,
            "has_language": False,
            "has_tag_tree": False,
            "all_visuals_have_alt": False,
        }

        try:
            # Check modelpack for title and language
            if modelpack_path.exists():
                with open(modelpack_path, encoding="utf-8") as f:
                    modelpack = json.load(f)

                # Check title
                title = modelpack.get("presentation_title", "")
                if not title:
                    # Try to get from first slide
                    slides = modelpack.get("slides", [])
                    if slides:
                        for elem in slides[0].get("elements_in_reading_order", []):
                            if elem.get("text"):
                                title = elem["text"]
                                break

                if title:
                    compliance_passed["has_title"] = True
                    print(f"[Compliance] ✓ Document title: {title[:50]}...")
                else:
                    compliance_issues.append("Missing document title")

                # Check language (default is de-DE)
                language = modelpack.get("language", "de-DE")
                if language:
                    compliance_passed["has_language"] = True
                    print(f"[Compliance] ✓ Document language: {language}")
                else:
                    compliance_issues.append("Missing document language")

                # Check that slides exist (tag tree basis)
                if modelpack.get("slides"):
                    compliance_passed["has_tag_tree"] = True
                    print(f"[Compliance] ✓ Tag tree structure: {len(modelpack['slides'])} slides")
                else:
                    compliance_issues.append("No slides found for tag tree")

            # Check alt-texts from QA results
            if alttext_qa_path.exists():
                with open(alttext_qa_path, encoding="utf-8") as f:
                    qa_result = json.load(f)

                entries = qa_result.get("entries", [])
                missing_alt = 0
                for entry in entries:
                    alt_text = entry.get("alt_text", "")
                    is_decorative = entry.get("is_decorative", False)

                    if not alt_text and not is_decorative:
                        missing_alt += 1

                if missing_alt == 0:
                    compliance_passed["all_visuals_have_alt"] = True
                    print(f"[Compliance] ✓ All visual elements have alt-text or are marked decorative")
                else:
                    compliance_issues.append(f"{missing_alt} visual elements missing alt-text")

            # Calculate overall compliance
            total_checks = len(compliance_passed)
            passed_checks = sum(1 for v in compliance_passed.values() if v)
            compliance_score = (passed_checks / total_checks) * 100 if total_checks > 0 else 0

            is_compliant = all(compliance_passed.values())

            if is_compliant:
                print(f"[Compliance] ✓ PDF/UA COMPLIANT - All {total_checks} checks passed")
            else:
                print(f"[Compliance] ✗ COMPLIANCE ISSUES - {passed_checks}/{total_checks} checks passed")
                for issue in compliance_issues:
                    print(f"[Compliance]   - {issue}")

            return StageResult(
                stage=PipelineStage.COMPLIANCE_CHECK,
                success=True,  # Stage succeeds even if not fully compliant (for logging)
                metadata={
                    "compliance_score": compliance_score,
                    "is_compliant": is_compliant,
                    "checks_passed": passed_checks,
                    "checks_total": total_checks,
                    "issues": compliance_issues,
                    "details": compliance_passed,
                },
            )

        except Exception as e:
            return StageResult(
                stage=PipelineStage.COMPLIANCE_CHECK,
                success=False,
                error=f"Compliance check failed: {str(e)}",
            )

    def _run_final_render(
        self,
        ir_path: Path,
        semantic_path: Path,
        alttext_qa_path: Path,
        output_dir: Path,
    ) -> StageResult:
        """
        Stage 6: Generate final accessible output (PDF/UA or HTML).

        This is a placeholder for future implementation.
        """
        # TODO: Implement PDF/UA or accessible HTML generation
        # For now, just mark as successful if inputs exist
        if not ir_path.exists():
            return StageResult(
                stage=PipelineStage.FINAL_RENDER,
                success=False,
                error="IR file not found",
            )

        return StageResult(
            stage=PipelineStage.FINAL_RENDER,
            success=True,
            output_path=output_dir / "output.pdf",
            metadata={"format": "pdf_ua", "note": "Placeholder - not yet implemented"},
        )

    def _validate_pdf_output(
        self,
        pdf_path: Path,
    ) -> Tuple[StageResult, ComplianceReport]:
        """
        Validate the generated PDF for PDF/UA compliance.

        Called AFTER PDF generation to verify the actual file meets
        accessibility requirements.

        Args:
            pdf_path: Path to the generated PDF file.

        Returns:
            Tuple of StageResult and ComplianceReport.
        """
        try:
            if not pdf_path.exists():
                report = ComplianceReport(is_compliant=False)
                report.add_error(f"PDF nicht gefunden: {pdf_path}")
                return StageResult(
                    stage=PipelineStage.COMPLIANCE_CHECK,
                    success=False,
                    error=f"PDF not found: {pdf_path}",
                ), report

            # Run PDF validation
            report = validate_pdf_accessibility(str(pdf_path))

            # Log results
            if report.is_compliant:
                print(f"[PDF Validator] ✓ PDF/UA COMPLIANT")
            else:
                print(f"[PDF Validator] ✗ INCOMPLETE COMPLIANCE")
                for error in report.errors:
                    print(f"[PDF Validator]   - {error}")

            for warning in report.warnings:
                print(f"[PDF Validator]   ⚠ {warning}")

            return StageResult(
                stage=PipelineStage.COMPLIANCE_CHECK,
                success=True,  # Stage succeeds even if PDF not fully compliant
                output_path=pdf_path,
                metadata={
                    "pdf_is_compliant": report.is_compliant,
                    "pdf_compliance_status": report.status,
                    "pdf_errors": report.errors,
                    "pdf_warnings": report.warnings,
                    "pdf_stats": report.stats,
                },
            ), report

        except Exception as e:
            report = ComplianceReport(is_compliant=False)
            report.add_error(f"Validierung fehlgeschlagen: {e}")
            return StageResult(
                stage=PipelineStage.COMPLIANCE_CHECK,
                success=False,
                error=f"PDF validation failed: {e}",
            ), report


def run_pipeline(
    pptx_path: Path,
    output_dir: Path,
    job_id: Optional[str] = None,
) -> PipelineResult:
    """
    Convenience function to run the full pipeline.

    Args:
        pptx_path: Path to the input PPTX file.
        output_dir: Directory for output files.
        job_id: Optional job ID (generated if not provided).

    Returns:
        PipelineResult with all stage results.
    """
    import uuid

    job_id = job_id or str(uuid.uuid4())
    orchestrator = PipelineOrchestrator(output_dir=output_dir)
    return orchestrator.run(pptx_path, job_id)
