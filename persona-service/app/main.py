"""
Persona-Check Service - Accessibility Persona Matcher
Analyzes documents against 7 realistic personas.
"""

import os
import uuid
import asyncio
from typing import Dict, Any, Optional
from datetime import datetime

from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

from app.config import PERSONAS, OLLAMA_CONFIG
from app.matcher import assess_text_for_personas
from app.parsers.pdf_parser import parse_pdf
from app.parsers.pptx_parser import parse_pptx
from app.parsers.docx_parser import parse_docx
from app.analyzers.readability import analyze_readability
from app.analyzers.structure import analyze_structure
from app.analyzers.contrast import analyze_colors
from app.analyzers.accessibility import analyze_accessibility
from app.analyzers.ai_analysis import analyze_with_ai
from app.analyzers.chapter_detection import detect_chapters
from app.analyzers.passive_detector import analyze_passive_voice
from app.report_generator import generate_report


app = FastAPI(
    title="Persona-Check Service",
    description="Accessibility analysis against 7 realistic personas",
    version="1.0.0"
)

cors_origins_env = os.getenv("CORS_ORIGINS", "")
cors_origins = [origin.strip() for origin in cors_origins_env.split(",") if origin.strip()]
if "*" in cors_origins:
    cors_origins = ["*"]

app.add_middleware(
    CORSMiddleware,
    allow_origins=cors_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# In-memory job storage (use Redis in production)
jobs: Dict[str, Dict[str, Any]] = {}


class JobStatus(BaseModel):
    job_id: str
    status: str
    progress: int
    current_stage: str
    created_at: str
    completed_at: Optional[str] = None
    error: Optional[str] = None


class AnalysisResponse(BaseModel):
    job_id: str
    status: str
    message: str


UPLOAD_DIR = "/app/uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)


def get_parser(filename: str):
    """Get the appropriate parser based on file extension."""
    ext = filename.lower().split(".")[-1]
    parsers = {
        "pdf": parse_pdf,
        "pptx": parse_pptx,
        "ppt": parse_pptx,
        "docx": parse_docx,
        "doc": parse_docx,
    }
    return parsers.get(ext)


async def run_analysis(job_id: str, file_path: str, filename: str):
    """Run the full analysis pipeline."""
    try:
        # Stage 1: Parse document
        jobs[job_id]["status"] = "parsing"
        jobs[job_id]["current_stage"] = "Dokument wird analysiert..."
        jobs[job_id]["progress"] = 10

        parser = get_parser(filename)
        if not parser:
            raise ValueError(f"Unsupported file type: {filename}")

        document_data = await parser(file_path)
        jobs[job_id]["progress"] = 25

        # Stage 2: Automatic analysis
        jobs[job_id]["status"] = "analyzing"
        jobs[job_id]["current_stage"] = "Automatische Tests..."
        jobs[job_id]["progress"] = 30

        analysis_results = {}

        # Readability analysis
        analysis_results["readability"] = analyze_readability(document_data["text"])
        jobs[job_id]["progress"] = 40

        # Structure analysis
        analysis_results["structure"] = analyze_structure(document_data)
        jobs[job_id]["progress"] = 50

        # Color/contrast analysis (if colors available)
        if document_data.get("colors"):
            analysis_results["colors"] = analyze_colors(document_data["colors"])
        jobs[job_id]["progress"] = 60

        # Accessibility checks
        analysis_results["accessibility"] = analyze_accessibility(document_data)
        jobs[job_id]["progress"] = 65

        # Passive voice analysis
        analysis_results["passive_voice"] = analyze_passive_voice(document_data["text"])
        jobs[job_id]["progress"] = 70

        # Stage 3: AI-powered analysis
        jobs[job_id]["status"] = "ai_analysis"
        jobs[job_id]["current_stage"] = "KI-Analyse (Tonalität, Fachbegriffe)..."
        jobs[job_id]["progress"] = 75

        try:
            analysis_results["ai"] = await analyze_with_ai(document_data["text"])
        except Exception as e:
            print(f"AI analysis failed: {e}")
            analysis_results["ai"] = {"error": str(e)}
        jobs[job_id]["progress"] = 90

        # Stage 4: Generate report
        jobs[job_id]["status"] = "generating_report"
        jobs[job_id]["current_stage"] = "Report wird erstellt..."

        report = generate_report(filename, document_data, analysis_results, PERSONAS)
        jobs[job_id]["progress"] = 100

        # Store results
        jobs[job_id]["status"] = "completed"
        jobs[job_id]["current_stage"] = "Fertig"
        jobs[job_id]["completed_at"] = datetime.utcnow().isoformat()
        jobs[job_id]["report"] = report

        # Cleanup file
        if os.path.exists(file_path):
            os.remove(file_path)

    except Exception as e:
        jobs[job_id]["status"] = "failed"
        jobs[job_id]["error"] = str(e)
        jobs[job_id]["current_stage"] = "Fehler"
        print(f"Analysis failed for job {job_id}: {e}")

        # Cleanup file
        if os.path.exists(file_path):
            os.remove(file_path)


@app.get("/health")
async def health_check():
    """Health check endpoint."""
    # Check Ollama connection
    ollama_available = False
    try:
        import ollama
        client = ollama.Client(host=OLLAMA_CONFIG["host"])
        client.list()
        ollama_available = True
    except Exception:
        pass

    return {
        "status": "ok",
        "personas": len(PERSONAS),
        "ollama": {
            "available": ollama_available,
            "host": OLLAMA_CONFIG["host"],
            "model": OLLAMA_CONFIG["model"],
        },
        "supported_formats": ["pdf", "pptx", "docx"],
    }


@app.post("/analyze", response_model=AnalysisResponse)
async def analyze_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...)
):
    """Upload and analyze a document against all personas."""
    # Validate file type
    if not file.filename:
        raise HTTPException(status_code=400, detail="Filename is required")

    ext = file.filename.lower().split(".")[-1]
    if ext not in ["pdf", "pptx", "ppt", "docx", "doc"]:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {ext}. Supported: PDF, PPTX, DOCX"
        )

    # Create job
    job_id = str(uuid.uuid4())
    file_path = os.path.join(UPLOAD_DIR, f"{job_id}_{file.filename}")

    # Save file
    content = await file.read()
    with open(file_path, "wb") as f:
        f.write(content)

    # Initialize job
    jobs[job_id] = {
        "job_id": job_id,
        "filename": file.filename,
        "status": "pending",
        "progress": 0,
        "current_stage": "In Warteschlange...",
        "created_at": datetime.utcnow().isoformat(),
        "completed_at": None,
        "error": None,
        "report": None,
    }

    # Start analysis in background
    background_tasks.add_task(run_analysis, job_id, file_path, file.filename)

    return AnalysisResponse(
        job_id=job_id,
        status="pending",
        message="Analyse gestartet"
    )


@app.get("/status/{job_id}")
async def get_job_status(job_id: str):
    """Get the status of an analysis job."""
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job nicht gefunden")

    job = jobs[job_id]
    return {
        "job_id": job_id,
        "status": job["status"],
        "progress": job["progress"],
        "current_stage": job["current_stage"],
        "created_at": job["created_at"],
        "completed_at": job.get("completed_at"),
        "error": job.get("error"),
    }


@app.get("/report/{job_id}")
async def get_report(job_id: str):
    """Get the full report for a completed job."""
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job nicht gefunden")

    job = jobs[job_id]
    if job["status"] != "completed":
        raise HTTPException(
            status_code=400,
            detail=f"Report nicht verfügbar. Status: {job['status']}"
        )

    return job["report"]


@app.get("/personas")
async def list_personas():
    """List all available personas."""
    return {
        "count": len(PERSONAS),
        "personas": [
            {
                "id": p.id,
                "name": p.name,
                "age": p.age,
                "description": p.description,
                "core_limitation": p.core_limitation,
                "icon": p.icon,
            }
            for p in PERSONAS.values()
        ]
    }


class ChapterRequest(BaseModel):
    """Request model for chapter detection."""
    srt: str


class ChapterResponse(BaseModel):
    """Response model for chapter detection."""
    chapters: list
    count: int


@app.post("/analyze-chapters", response_model=ChapterResponse)
async def analyze_chapters(request: ChapterRequest):
    """
    Analyze SRT subtitles and detect logical chapter breaks.

    Uses LLM to identify transitions between speakers, topics, or sections
    and generates appropriate chapter titles.
    """
    if not request.srt or len(request.srt.strip()) < 10:
        raise HTTPException(
            status_code=400,
            detail="SRT content is required and must not be empty"
        )

    try:
        chapters = await detect_chapters(request.srt)
        return ChapterResponse(
            chapters=chapters,
            count=len(chapters)
        )
    except Exception as e:
        print(f"[ChapterAnalysis] Error: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Chapter detection failed: {str(e)}"
        )


# ===========================================
# TEXT PERSONA ASSESSMENT (for Einfache Sprache)
# ===========================================

class TextStats(BaseModel):
    """Text statistics from frontend analysis."""
    characters: int = 0
    words: int = 0
    sentences: int = 0
    avgWordsPerSentence: float = 0
    complexWords: int = 0
    passiveConstructions: int = 0
    genitiveConstructions: int = 0
    abbreviations: int = 0
    replaceableWords: int = 0
    flesch_score: float = 50  # Optional Flesch reading ease score


class TextIssue(BaseModel):
    """Single text issue from analysis."""
    type: str
    text: str
    suggestion: Optional[str] = None


class TextAssessmentRequest(BaseModel):
    """Request model for text persona assessment."""
    stats: TextStats
    issues: list[TextIssue] = []


class PersonaAssessment(BaseModel):
    """Assessment result for a single persona."""
    personaId: str
    status: str  # 'accessible' | 'limited' | 'blocked'
    matchScore: int  # 0-100
    specificIssues: list[str]


class TextAssessmentResponse(BaseModel):
    """Response model for text persona assessment."""
    personaChecks: list[PersonaAssessment]


@app.post("/assess-text", response_model=TextAssessmentResponse)
async def assess_text(request: TextAssessmentRequest):
    """
    Assess text accessibility for all 7 personas.

    Takes text statistics and detected issues from the frontend text analysis
    and returns an accessibility assessment for each persona.
    """
    try:
        # Convert Pydantic models to dicts for the matcher
        stats_dict = request.stats.model_dump()
        issues_list = [issue.model_dump() for issue in request.issues]

        # Run persona assessment
        assessments = assess_text_for_personas(stats_dict, issues_list)

        return TextAssessmentResponse(personaChecks=assessments)

    except Exception as e:
        print(f"[TextAssessment] Error: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Text assessment failed: {str(e)}"
        )


@app.get("/personas-metadata")
async def get_personas_metadata():
    """
    Get metadata for all personas (for frontend display).

    Returns icons, names, descriptions, and limitations.
    """
    return {
        "personas": {
            p_id: {
                "id": p.id,
                "name": p.name,
                "age": p.age,
                "description": p.description,
                "core_limitation": p.core_limitation,
                "icon": p.icon,
            }
            for p_id, p in PERSONAS.items()
        }
    }


# ===========================================
# CUSTOM PERSONAS
# ===========================================

# In-memory storage for custom personas (use database in production)
custom_personas: Dict[str, Dict[str, Any]] = {}

# Available tests that users can select from
AVAILABLE_TESTS = [
    {"id": "contrast", "name": "Kontrast", "category": "visual", "description": "WCAG AA Kontrastverhältnis (4.5:1)"},
    {"id": "font_size", "name": "Schriftgröße", "category": "visual", "description": "Mindestschriftgröße 12pt"},
    {"id": "alt_texts", "name": "Alternativtexte", "category": "visual", "description": "Bilder mit Beschreibung"},
    {"id": "readability", "name": "Lesbarkeit", "category": "language", "description": "Flesch-Index für Textverständlichkeit"},
    {"id": "sentence_length", "name": "Satzlänge", "category": "language", "description": "Durchschnittliche Wörter pro Satz"},
    {"id": "passive_voice", "name": "Passiv-Sätze", "category": "language", "description": "Anteil Passiv-Konstruktionen"},
    {"id": "jargon", "name": "Fachbegriffe", "category": "language", "description": "Unbekannte/technische Begriffe"},
    {"id": "structure", "name": "Struktur", "category": "structure", "description": "Überschriften und Gliederung"},
    {"id": "paragraph_length", "name": "Absatzlänge", "category": "structure", "description": "Wörter pro Absatz"},
    {"id": "tonality", "name": "Tonalität", "category": "emotional", "description": "Freundlicher vs. bedrohlicher Ton"},
    {"id": "color_only", "name": "Nur Farbe", "category": "visual", "description": "Information nur durch Farbe vermittelt"},
    {"id": "simple_language", "name": "Leichte Sprache", "category": "language", "description": "Entspricht Regeln für Leichte Sprache"},
]


class CustomPersonaCreate(BaseModel):
    """Request model for creating a custom persona."""
    name: str
    age: int = 40
    description: str
    core_limitation: str
    icon: str = "👤"
    relevant_tests: list[str]


class CustomPersonaResponse(BaseModel):
    """Response model for custom persona."""
    id: str
    name: str
    age: int
    description: str
    core_limitation: str
    icon: str
    relevant_tests: list[str]


@app.get("/available-tests")
async def get_available_tests():
    """
    Get list of available tests that can be assigned to custom personas.
    """
    return {
        "tests": AVAILABLE_TESTS,
        "categories": {
            "visual": "Visuell",
            "language": "Sprache",
            "structure": "Struktur",
            "emotional": "Emotional/Kognitiv",
        }
    }


@app.post("/custom-personas", response_model=CustomPersonaResponse)
async def create_custom_persona(persona: CustomPersonaCreate):
    """
    Create a custom persona for document analysis.

    The persona will be included in future analyses alongside the 7 default personas.
    """
    # Generate ID from name
    persona_id = f"custom_{persona.name.lower().replace(' ', '_')}"

    # Validate tests
    valid_test_ids = {t["id"] for t in AVAILABLE_TESTS}
    invalid_tests = [t for t in persona.relevant_tests if t not in valid_test_ids]
    if invalid_tests:
        raise HTTPException(
            status_code=400,
            detail=f"Ungültige Tests: {invalid_tests}. Verfügbar: {list(valid_test_ids)}"
        )

    # Store custom persona
    custom_personas[persona_id] = {
        "id": persona_id,
        "name": persona.name,
        "age": persona.age,
        "description": persona.description,
        "core_limitation": persona.core_limitation,
        "icon": persona.icon,
        "relevant_tests": persona.relevant_tests,
        "is_custom": True,
    }

    return CustomPersonaResponse(
        id=persona_id,
        name=persona.name,
        age=persona.age,
        description=persona.description,
        core_limitation=persona.core_limitation,
        icon=persona.icon,
        relevant_tests=persona.relevant_tests,
    )


@app.get("/custom-personas")
async def list_custom_personas():
    """
    List all custom personas.
    """
    return {
        "count": len(custom_personas),
        "personas": list(custom_personas.values()),
    }


@app.delete("/custom-personas/{persona_id}")
async def delete_custom_persona(persona_id: str):
    """
    Delete a custom persona.
    """
    if persona_id not in custom_personas:
        raise HTTPException(status_code=404, detail="Persona nicht gefunden")

    del custom_personas[persona_id]
    return {"message": "Persona gelöscht", "id": persona_id}


# ===========================================
# TEXT OPTIMIZATION FOR PERSONAS
# ===========================================

# Persona-specific optimization prompts
PERSONA_OPTIMIZATION_PROMPTS = {
    "ingrid": {
        "name": "Ingrid (78)",
        "focus": "Sehbeeinträchtigung, Digital Immigrant",
        "prompt": """Du schreibst für Ingrid, 78 Jahre alt, mit altersbedingter Makuladegeneration.

WICHTIGE REGELN:
- Kurze, klare Sätze (max. 15 Wörter)
- Aktive Formulierungen statt Passiv
- Keine Abkürzungen
- Einfache Wörter, keine Fremdwörter
- Wichtige Informationen am Satzanfang
- Klare Struktur mit Absätzen

Schreibe den folgenden Text so um, dass Ingrid ihn leicht verstehen kann:"""
    },
    "mehmet": {
        "name": "Mehmet (34)",
        "focus": "Deutsch B1, Fluchthintergrund",
        "prompt": """Du schreibst für Mehmet, 34 Jahre alt, mit Deutsch auf B1-Niveau und Fluchthintergrund.

WICHTIGE REGELN:
- Einfaches Deutsch (B1-Niveau)
- Kurze Sätze (max. 12-15 Wörter)
- KEIN Passiv - nur Aktiv-Sätze
- Keine Fachbegriffe ohne Erklärung
- Keine Redewendungen oder Sprichwörter
- Freundlicher, nicht bedrohlicher Ton
- Bei Behörden-Kontext: Erkläre was zu tun ist

Schreibe den folgenden Text so um, dass Mehmet ihn verstehen kann:"""
    },
    "fatima": {
        "name": "Fatima (29)",
        "focus": "Analphabetin, A2 mündlich",
        "prompt": """Du schreibst für Fatima, 29 Jahre alt, die nicht lesen kann und nur A2-Deutsch spricht.

WICHTIGE REGELN:
- Extrem einfache Sprache (Leichte Sprache Niveau)
- Sehr kurze Sätze (max. 8-10 Wörter)
- Ein Gedanke pro Satz
- Nur Grundwortschatz
- Wiederholungen sind OK
- Schritt-für-Schritt Anleitungen
- Zahlen als Ziffern schreiben

Schreibe den folgenden Text in Leichter Sprache:"""
    },
    "thomas": {
        "name": "Thomas (42)",
        "focus": "ADHS & Legasthenie",
        "prompt": """Du schreibst für Thomas, 42 Jahre alt, mit ADHS und Legasthenie.

WICHTIGE REGELN:
- Klare visuelle Struktur mit Überschriften
- Kurze Absätze (max. 3 Sätze)
- Aufzählungen und Listen verwenden
- Wichtiges fett markieren mit **text**
- Keine langen Textblöcke
- Handlungsaufforderungen klar formulieren
- Ablenkungen minimieren

Schreibe den folgenden Text strukturiert für Thomas um:"""
    },
    "aisha": {
        "name": "Aisha (67)",
        "focus": "Gehörlos, DGS-Muttersprachlerin",
        "prompt": """Du schreibst für Aisha, 67 Jahre alt, gehörlos mit Deutscher Gebärdensprache als Muttersprache.

WICHTIGE REGELN:
- KEIN Passiv - NUR Aktiv-Sätze (Subjekt-Verb-Objekt)
- Kurze, direkte Sätze
- Wer macht was? Immer klar benennen
- Keine verschachtelten Nebensätze
- Zeitliche Abfolge klar machen
- Konkrete statt abstrakte Begriffe

Schreibe den folgenden Text ohne Passiv für Aisha um:"""
    },
    "sandra": {
        "name": "Sandra (35)",
        "focus": "Alleinerziehend, Zeitdruck",
        "prompt": """Du schreibst für Sandra, 35 Jahre alt, alleinerziehend und unter chronischem Zeitdruck.

WICHTIGE REGELN:
- Wichtigstes zuerst (Pyramiden-Prinzip)
- Klare Handlungsanweisungen
- Fristen und Termine hervorheben
- Unnötige Details weglassen
- Was muss Sandra TUN? Klar benennen
- Zeitschätzungen geben wenn möglich

Schreibe den folgenden Text kurz und handlungsorientiert für Sandra um:"""
    },
    "viktor": {
        "name": "Viktor (52)",
        "focus": "Farbenblindheit, Depression",
        "prompt": """Du schreibst für Viktor, 52 Jahre alt, mit Depression und Rot-Grün-Schwäche.

WICHTIGE REGELN:
- Positiver, ermutigender Ton
- Keine bedrohlichen Formulierungen
- Kleine, machbare Schritte aufzeigen
- Erfolge und Fortschritte betonen
- "Sie müssen" vermeiden, besser "Sie können"
- Hilfsangebote erwähnen
- Nicht überwältigend

Schreibe den folgenden Text ermutigend und nicht bedrohlich für Viktor um:"""
    },
}


class OptimizeTextRequest(BaseModel):
    """Request model for persona-specific text optimization."""
    text: str
    persona_id: str


class OptimizeTextResponse(BaseModel):
    """Response model for optimized text."""
    persona_id: str
    persona_name: str
    original_text: str
    optimized_text: str
    optimization_notes: list[str]


@app.post("/optimize-text", response_model=OptimizeTextResponse)
async def optimize_text_for_persona(request: OptimizeTextRequest):
    """
    Optimize text for a specific persona using AI.

    Takes the original text and rewrites it according to the
    accessibility needs of the specified persona.
    """
    persona_id = request.persona_id
    text = request.text

    if not text or len(text.strip()) < 10:
        raise HTTPException(
            status_code=400,
            detail="Text muss mindestens 10 Zeichen lang sein"
        )

    # Get persona prompt
    if persona_id not in PERSONA_OPTIMIZATION_PROMPTS:
        raise HTTPException(
            status_code=404,
            detail=f"Persona '{persona_id}' nicht gefunden. Verfügbar: {list(PERSONA_OPTIMIZATION_PROMPTS.keys())}"
        )

    persona_config = PERSONA_OPTIMIZATION_PROMPTS[persona_id]

    try:
        import ollama

        client = ollama.Client(host=OLLAMA_CONFIG["host"])

        # Build the prompt
        full_prompt = f"""{persona_config["prompt"]}

---
ORIGINALTEXT:
{text}
---

Gib NUR den umgeschriebenen Text zurück, ohne Erklärungen oder Kommentare."""

        # Call Ollama
        response = client.generate(
            model=OLLAMA_CONFIG["model"],
            prompt=full_prompt,
            options={
                "temperature": 0.3,  # Lower temperature for more consistent output
                "num_predict": 2000,
            }
        )

        optimized_text = response.get("response", "").strip()

        if not optimized_text:
            raise HTTPException(
                status_code=500,
                detail="Keine Antwort vom KI-Modell erhalten"
            )

        # Generate optimization notes based on persona
        notes = _generate_optimization_notes(persona_id)

        return OptimizeTextResponse(
            persona_id=persona_id,
            persona_name=persona_config["name"],
            original_text=text,
            optimized_text=optimized_text,
            optimization_notes=notes,
        )

    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="Ollama-Bibliothek nicht installiert"
        )
    except Exception as e:
        print(f"[OptimizeText] Error for {persona_id}: {e}")
        raise HTTPException(
            status_code=503,
            detail=f"KI-Service nicht verfügbar: {str(e)}"
        )


def _generate_optimization_notes(persona_id: str) -> list[str]:
    """Generate notes about what was optimized for the persona."""
    notes_map = {
        "ingrid": [
            "Sätze gekürzt für bessere Lesbarkeit",
            "Passiv durch Aktiv ersetzt",
            "Fremdwörter vereinfacht",
        ],
        "mehmet": [
            "Auf B1-Deutsch-Niveau vereinfacht",
            "Passiv-Konstruktionen entfernt",
            "Fachbegriffe erklärt oder ersetzt",
        ],
        "fatima": [
            "In Leichte Sprache übersetzt",
            "Sehr kurze Sätze verwendet",
            "Nur Grundwortschatz verwendet",
        ],
        "thomas": [
            "Klare Struktur mit Überschriften",
            "Wichtiges hervorgehoben",
            "Absätze verkürzt",
        ],
        "aisha": [
            "Alle Passiv-Sätze in Aktiv umgewandelt",
            "Subjekt-Verb-Objekt Struktur",
            "Keine Nebensätze",
        ],
        "sandra": [
            "Wichtigstes zuerst",
            "Handlungsanweisungen klar formuliert",
            "Unnötige Details entfernt",
        ],
        "viktor": [
            "Bedrohliche Formulierungen entfernt",
            "Positiver, ermutigender Ton",
            "Kleine Schritte aufgezeigt",
        ],
    }
    return notes_map.get(persona_id, ["Text optimiert"])


# ===========================================
# FIX CRITICAL ISSUES IN PPTX
# ===========================================

from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Pt
import tempfile
import re


class FixCriticalResponse(BaseModel):
    """Response model for fix-critical endpoint."""
    persona_id: str
    persona_name: str
    slides_modified: int
    changes: list[dict]


def _identify_critical_text(text: str, persona_id: str) -> bool:
    """
    Determine if a text passage likely has critical issues for the persona.
    Uses heuristics to identify problematic text without full AI analysis.
    """
    if len(text.strip()) < 20:
        return False

    # Count sentences
    sentences = re.split(r'[.!?]+', text)
    sentences = [s.strip() for s in sentences if s.strip()]

    # Average words per sentence
    if sentences:
        avg_words = sum(len(s.split()) for s in sentences) / len(sentences)
    else:
        avg_words = 0

    # Check for passive voice indicators (German)
    passive_indicators = ['wird', 'werden', 'wurde', 'wurden', 'worden', 'geworden']
    has_passive = any(indicator in text.lower() for indicator in passive_indicators)

    # Check for complex words (long words often indicate complexity)
    words = text.split()
    long_words = sum(1 for w in words if len(w) > 12)
    long_word_ratio = long_words / len(words) if words else 0

    # Persona-specific critical thresholds
    if persona_id == "ingrid":
        # Needs simple language, no passive
        return avg_words > 15 or has_passive or long_word_ratio > 0.1

    elif persona_id == "mehmet":
        # B1 level, no passive, no jargon
        return avg_words > 15 or has_passive or long_word_ratio > 0.15

    elif persona_id == "fatima":
        # Very simple language needed
        return avg_words > 10 or long_word_ratio > 0.05

    elif persona_id == "thomas":
        # Needs structure, short paragraphs
        return len(text) > 300 or avg_words > 20

    elif persona_id == "aisha":
        # No passive voice, simple grammar
        return has_passive or avg_words > 12

    elif persona_id == "sandra":
        # Concise, action-oriented
        return len(text) > 200 or avg_words > 18

    elif persona_id == "viktor":
        # Positive tone, not overwhelming
        negative_words = ['müssen', 'pflicht', 'zwingend', 'frist', 'strafe', 'versäumnis']
        has_negative = any(word in text.lower() for word in negative_words)
        return has_negative or len(text) > 250

    return False


def _optimize_text_sync(text: str, persona_id: str) -> str:
    """
    Synchronously optimize text for a persona using Ollama.
    """
    if persona_id not in PERSONA_OPTIMIZATION_PROMPTS:
        return text

    persona_config = PERSONA_OPTIMIZATION_PROMPTS[persona_id]

    try:
        import ollama
        client = ollama.Client(host=OLLAMA_CONFIG["host"])

        full_prompt = f"""{persona_config["prompt"]}

---
ORIGINALTEXT:
{text}
---

Gib NUR den umgeschriebenen Text zurück, ohne Erklärungen, Kommentare oder Anführungszeichen."""

        response = client.generate(
            model=OLLAMA_CONFIG["model"],
            prompt=full_prompt,
            options={
                "temperature": 0.3,
                "num_predict": 2000,
            }
        )

        optimized = response.get("response", "").strip()

        # Clean up any quotes or markers the LLM might have added
        if optimized.startswith('"') and optimized.endswith('"'):
            optimized = optimized[1:-1]
        if optimized.startswith("'") and optimized.endswith("'"):
            optimized = optimized[1:-1]

        return optimized if optimized else text

    except Exception as e:
        print(f"[FixCritical] Optimization failed for text: {e}")
        return text


@app.post("/fix-critical")
async def fix_critical_issues(
    file: UploadFile = File(...),
    persona_id: str = Form(...),
    job_id: Optional[str] = Form(None),
):
    """
    Fix critical accessibility issues in a PPTX file for a specific persona.

    Takes a PPTX file, identifies text passages with critical issues for the
    specified persona, rewrites them using AI, and returns the modified PPTX.
    """
    # Validate persona
    if persona_id not in PERSONA_OPTIMIZATION_PROMPTS:
        raise HTTPException(
            status_code=400,
            detail=f"Ungültige Persona: {persona_id}. Verfügbar: {list(PERSONA_OPTIMIZATION_PROMPTS.keys())}"
        )

    # Validate file type
    if not file.filename or not file.filename.lower().endswith(('.pptx', '.ppt')):
        raise HTTPException(
            status_code=400,
            detail="Nur PowerPoint-Dateien (.pptx) werden unterstützt"
        )

    # Save uploaded file temporarily
    temp_input = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    temp_output = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")

    try:
        content = await file.read()
        temp_input.write(content)
        temp_input.close()

        # Open the presentation
        prs = Presentation(temp_input.name)

        slides_modified = 0
        changes = []

        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_changed = False

            # Process all shapes with text
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    # Collect full paragraph text
                    para_text = "".join(run.text for run in paragraph.runs)

                    if not para_text.strip():
                        continue

                    # Check if this text needs fixing
                    if not _identify_critical_text(para_text, persona_id):
                        continue

                    # Optimize the text
                    optimized_text = _optimize_text_sync(para_text, persona_id)

                    if optimized_text != para_text:
                        # Replace text in runs while preserving formatting
                        # Strategy: Put all text in first run, clear others
                        if paragraph.runs:
                            first_run = paragraph.runs[0]
                            first_run.text = optimized_text

                            # Clear remaining runs
                            for run in paragraph.runs[1:]:
                                run.text = ""

                            slide_changed = True
                            changes.append({
                                "slide": slide_num,
                                "original": para_text[:100] + "..." if len(para_text) > 100 else para_text,
                                "optimized": optimized_text[:100] + "..." if len(optimized_text) > 100 else optimized_text,
                            })

            if slide_changed:
                slides_modified += 1

        # Save modified presentation
        temp_output.close()
        prs.save(temp_output.name)

        print(f"[FixCritical] Modified {slides_modified} slides with {len(changes)} text changes for persona {persona_id}")

        # Return the file
        return FileResponse(
            path=temp_output.name,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=file.filename.replace(".pptx", f"_{persona_id}_korrigiert.pptx"),
            background=None,  # Don't delete file until response is sent
        )

    except Exception as e:
        print(f"[FixCritical] Error: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"PPTX-Korrektur fehlgeschlagen: {str(e)}"
        )

    finally:
        # Cleanup input file
        try:
            os.remove(temp_input.name)
        except Exception:
            pass
        # Note: output file will be cleaned up by FileResponse or garbage collection
