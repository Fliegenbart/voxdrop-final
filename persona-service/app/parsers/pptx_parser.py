"""
PPTX Parser - Extract text, structure, and metadata from PowerPoint files.
Uses python-pptx for parsing.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from typing import Dict, Any, List
import os


async def parse_pptx(file_path: str) -> Dict[str, Any]:
    """
    Parse a PPTX file and extract relevant content.

    Returns:
        dict with:
        - text: Full text content
        - slides: List of slide contents
        - headings: Detected headings (slide titles)
        - images: Image information
        - colors: Detected colors
        - fonts: Font information
        - metadata: Document metadata
    """
    prs = Presentation(file_path)

    full_text = ""
    slides = []
    headings = []
    images = []
    colors = set()
    fonts = set()

    for slide_num, slide in enumerate(prs.slides):
        slide_text = ""
        slide_content = {
            "slide_number": slide_num + 1,
            "title": "",
            "text": "",
            "shapes": [],
            "has_notes": False,
        }

        # Get slide title
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            slide_content["title"] = title_text
            if title_text:
                headings.append({
                    "text": title_text,
                    "level": 1,
                    "slide": slide_num + 1,
                })

        # Process all shapes
        for shape in slide.shapes:
            # Handle text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ""
                    for run in paragraph.runs:
                        para_text += run.text

                        # Collect font info
                        if run.font:
                            font_name = run.font.name or "Default"
                            font_size = run.font.size.pt if run.font.size else 12
                            fonts.add((font_name, font_size))

                            # Collect text colors (handle _NoneColor gracefully)
                            try:
                                if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                    rgb = run.font.color.rgb
                                    colors.add((rgb.red, rgb.green, rgb.blue))
                            except (AttributeError, TypeError):
                                pass  # Skip colors that can't be extracted

                    slide_text += para_text + "\n"

            # Handle images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({
                    "slide": slide_num + 1,
                    "has_alt_text": bool(shape.name and not shape.name.startswith("Picture")),
                    "name": shape.name,
                })

            # Handle tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    slide_text += " | ".join(row_text) + "\n"

        # Check for speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_content["has_notes"] = True
                slide_content["notes"] = notes_text

        slide_content["text"] = slide_text.strip()
        full_text += slide_text + "\n"
        slides.append(slide_content)

    # Extract colors from slide backgrounds and shapes
    for slide in prs.slides:
        try:
            if slide.background.fill.solid():
                bg_color = slide.background.fill.fore_color
                if hasattr(bg_color, 'rgb') and bg_color.rgb:
                    rgb = bg_color.rgb
                    colors.add((rgb.red, rgb.green, rgb.blue))
        except Exception:
            pass

    # Get core properties (metadata)
    metadata = {}
    try:
        core_props = prs.core_properties
        metadata = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
        }
    except Exception:
        metadata = {"title": "", "author": "", "subject": ""}

    return {
        "text": full_text.strip(),
        "slides": slides,
        "slide_count": len(slides),
        "headings": headings,
        "images": images,
        "image_count": len(images),
        "colors": list(colors),
        "fonts": list(fonts),
        "metadata": metadata,
        "format": "pptx",
    }
