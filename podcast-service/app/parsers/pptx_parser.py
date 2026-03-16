"""
PPTX Parser - Extracts text and images from PowerPoint files

Enhanced with VLM image descriptions for podcast generation.
"""
import os
import io
import base64
import asyncio
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from app.config import VLM_CONFIG


def extract_slide_content(pptx_path: str) -> List[Dict[str, Any]]:
    """
    Extract text and images from each slide of a PPTX file.

    Returns:
        List of slide data with text and images
    """
    prs = Presentation(pptx_path)
    slides_data = []
    slide_w = int(prs.slide_width) if prs.slide_width else 0
    slide_h = int(prs.slide_height) if prs.slide_height else 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": slide_idx,
            "title": "",
            "content": [],
            "speaker_notes": "",
            "images": [],
            "slide_width": slide_w,
            "slide_height": slide_h,
        }

        # Extract shapes
        for shape in slide.shapes:
            # Title
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    slide_data["title"] = shape.text.strip()
                else:
                    text = shape.text.strip()
                    if text:
                        slide_data["content"].append(text)

            # Images
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    image_bytes = image.blob

                    # Convert to base64 for potential AI description
                    img_base64 = base64.b64encode(image_bytes).decode('utf-8')

                    slide_data["images"].append({
                        "data": img_base64,
                        "content_type": image.content_type,
                        "width": int(shape.width) if shape.width else 0,
                        "height": int(shape.height) if shape.height else 0,
                    })
                except Exception as e:
                    print(f"[PPTX Parser] Error extracting image from slide {slide_idx}: {e}")

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                slide_data["speaker_notes"] = notes_frame.text.strip()

        slides_data.append(slide_data)

    return slides_data


def get_presentation_metadata(pptx_path: str) -> Dict[str, Any]:
    """
    Get metadata from a PPTX file.
    """
    prs = Presentation(pptx_path)

    return {
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
    }


def slides_to_text(slides_data: List[Dict[str, Any]]) -> str:
    """
    Convert slides data to a structured text representation.
    """
    text_parts = []

    for slide in slides_data:
        text_parts.append(f"=== Folie {slide['slide_number']} ===")

        if slide["title"]:
            text_parts.append(f"Titel: {slide['title']}")

        if slide["content"]:
            text_parts.append("Inhalt:")
            for content in slide["content"]:
                text_parts.append(f"  - {content}")

        if slide["speaker_notes"]:
            text_parts.append(f"Notizen: {slide['speaker_notes']}")

        if slide["images"]:
            # Include image descriptions if available
            for i, img in enumerate(slide["images"]):
                if img.get("description"):
                    text_parts.append(f"Bild {i+1}: {img['description']}")
                else:
                    text_parts.append(f"[Bild {i+1}: Keine Beschreibung verfuegbar]")

        text_parts.append("")

    return "\n".join(text_parts)


async def extract_slides_with_descriptions(
    pptx_path: str,
    progress_callback=None
) -> List[Dict[str, Any]]:
    """
    Extract slides and add VLM-generated image descriptions.

    This is the enhanced extraction function that uses the VLM captioner
    to generate spoken-friendly descriptions for all images.

    Args:
        pptx_path: Path to the PPTX file
        progress_callback: Optional callback(percent, stage)

    Returns:
        List of slide data with image descriptions
    """
    # First extract slides in a worker thread to avoid blocking the event loop
    slides_data = await asyncio.to_thread(extract_slide_content, pptx_path)

    if not slides_data:
        return slides_data

    # Check if VLM is enabled
    if not VLM_CONFIG.get("enabled", True):
        print("[PPTX Parser] VLM disabled, skipping image descriptions")
        return slides_data

    # Count images
    total_images = sum(len(s.get("images", [])) for s in slides_data)
    if total_images == 0:
        print("[PPTX Parser] No images found in presentation")
        return slides_data

    print(f"[PPTX Parser] Found {total_images} images, starting VLM description...")

    # Import here to avoid circular imports
    from app.processors.vlm_captioner import describe_slide_images

    # Add descriptions
    slides_data = await describe_slide_images(slides_data, progress_callback)

    return slides_data
