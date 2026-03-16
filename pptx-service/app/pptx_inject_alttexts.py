#!/usr/bin/env python3
"""
pptx_inject_alttexts.py
Injects alt-texts into a PPTX file for accessibility.

Takes the original PPTX and a dict of shape_id → alt_text mappings,
then sets the description (alt-text) attribute on each shape.

Usage:
  python pptx_inject_alttexts.py input.pptx alttexts.json -o enhanced.pptx
"""

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.util import Inches


def load_json(path: str) -> Dict[str, Any]:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def get_all_shapes(slide) -> List:
    """Recursively get all shapes including grouped shapes."""
    shapes = []
    for shape in slide.shapes:
        shapes.append(shape)
        # Handle grouped shapes
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                shapes.append(child)
    return shapes


def inject_alttexts(pptx_path: str, alttexts: Dict[str, str], output_path: str) -> Dict[str, Any]:
    """
    Inject alt-texts into a PPTX file.

    Args:
        pptx_path: Path to the original PPTX file
        alttexts: Dict mapping shape_id (as string) → alt-text
        output_path: Path to save the enhanced PPTX

    Returns:
        Dict with statistics about the injection
    """
    prs = Presentation(pptx_path)

    stats = {
        "total_shapes": 0,
        "shapes_with_alttext": 0,
        "alttexts_applied": 0,
        "slides_processed": 0,
    }

    for slide_idx, slide in enumerate(prs.slides):
        stats["slides_processed"] += 1
        shapes = get_all_shapes(slide)

        for shape in shapes:
            stats["total_shapes"] += 1

            # Get shape_id - python-pptx uses shape_id attribute
            shape_id = str(shape.shape_id) if hasattr(shape, "shape_id") else None

            if not shape_id:
                continue

            # Check if we have an alt-text for this shape
            alttext = alttexts.get(shape_id)

            if alttext:
                try:
                    # Set the alt-text via the _element's descr attribute
                    # This is the accessibility description in OOXML
                    nvSpPr = shape._element.find(
                        ".//{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr"
                    )
                    if nvSpPr is None:
                        nvSpPr = shape._element.find(
                            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr"
                        )
                    if nvSpPr is None:
                        nvSpPr = shape._element.find(
                            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}nvGrpSpPr"
                        )

                    if nvSpPr is not None:
                        cNvPr = nvSpPr.find(
                            "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
                        )
                        if cNvPr is None:
                            cNvPr = nvSpPr.find(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr"
                            )

                        if cNvPr is not None:
                            cNvPr.set("descr", alttext)
                            stats["alttexts_applied"] += 1
                        else:
                            # Alternative: try to set on the element directly
                            # This works for some shape types
                            try:
                                shape._element.set("descr", alttext)
                                stats["alttexts_applied"] += 1
                            except Exception:
                                pass
                    else:
                        # Fallback: try setting on nvCxnSpPr for connectors
                        # or directly on the shape element
                        try:
                            # For pictures and other shapes, try direct approach
                            for elem in shape._element.iter():
                                if elem.tag.endswith("cNvPr"):
                                    elem.set("descr", alttext)
                                    stats["alttexts_applied"] += 1
                                    break
                        except Exception:
                            pass

                except Exception as e:
                    print(f"Warning: Could not set alt-text for shape {shape_id}: {e}")

            # Check if shape already has alt-text
            try:
                existing_descr = None
                for elem in shape._element.iter():
                    if elem.tag.endswith("cNvPr"):
                        existing_descr = elem.get("descr")
                        break
                if existing_descr:
                    stats["shapes_with_alttext"] += 1
            except Exception:
                pass

    # Save the enhanced presentation
    prs.save(output_path)

    return stats


def main():
    ap = argparse.ArgumentParser(description="Inject alt-texts into a PPTX file")
    ap.add_argument("pptx_path", help="Path to the original PPTX file")
    ap.add_argument("alttexts_json", help="Path to alttexts.json (or JSON with 'alttexts' key)")
    ap.add_argument("-o", "--out", required=True, help="Output path for enhanced PPTX")

    args = ap.parse_args()

    # Load alt-texts
    data = load_json(args.alttexts_json)
    alttexts = data.get("alttexts", data)  # Support both {alttexts: {...}} and direct dict

    # Inject and save
    stats = inject_alttexts(args.pptx_path, alttexts, args.out)

    print(f"Processed {stats['slides_processed']} slides")
    print(f"Total shapes: {stats['total_shapes']}")
    print(f"Alt-texts applied: {stats['alttexts_applied']}")
    print(f"Wrote: {args.out}")


if __name__ == "__main__":
    main()
