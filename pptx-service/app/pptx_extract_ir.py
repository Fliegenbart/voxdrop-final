#!/usr/bin/env python3
"""
pptx_extract_ir.py
Deterministic PPTX → JSON IR extractor (no AI).

Usage:
  pip install python-pptx
  python pptx_extract_ir.py input.pptx -o output.json

Notes:
- Reads text, geometry, grouping, tables, images, basic chart info, speaker notes.
- Does NOT try to infer meaning/reading order; it only extracts facts.
"""

import argparse
import hashlib
import json
import os
from datetime import datetime
from typing import Any, Dict, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement  # type: ignore


EMU_PER_INCH = 914400
POINTS_PER_INCH = 72


def emu_to_points(emu: int) -> float:
    return round((emu / EMU_PER_INCH) * POINTS_PER_INCH, 2)


def safe_int(x: Any, default: int = 0) -> int:
    try:
        return int(x)
    except Exception:
        return default


def sha256_bytes(data: bytes) -> str:
    h = hashlib.sha256()
    h.update(data)
    return h.hexdigest()


def get_slide_size_points(prs: Presentation) -> Dict[str, float]:
    return {
        "width_pt": emu_to_points(prs.slide_width),
        "height_pt": emu_to_points(prs.slide_height),
    }


def extract_textframe(shape) -> Optional[Dict[str, Any]]:
    """Extracts paragraph/run structure deterministically."""
    if not getattr(shape, "has_text_frame", False):
        return None

    tf = shape.text_frame
    paragraphs = []
    full_text_parts = []

    for p_idx, p in enumerate(tf.paragraphs):
        runs = []
        p_text_parts = []
        for r_idx, r in enumerate(p.runs):
            t = r.text or ""
            p_text_parts.append(t)
            runs.append(
                {
                    "run_index": r_idx,
                    "text": t,
                    "bold": bool(r.font.bold) if r.font is not None else None,
                    "italic": bool(r.font.italic) if r.font is not None else None,
                    "underline": bool(r.font.underline) if r.font is not None else None,
                    "font_name": r.font.name if r.font is not None else None,
                    "font_size_pt": float(r.font.size.pt) if (r.font is not None and r.font.size) else None,
                }
            )

        p_text = "".join(p_text_parts)
        full_text_parts.append(p_text)
        paragraphs.append(
            {
                "paragraph_index": p_idx,
                "text": p_text,
                "level": safe_int(getattr(p, "level", 0), 0),
                "alignment": str(getattr(p, "alignment", None)),
                "runs": runs,
            }
        )

    full_text = "\n".join(full_text_parts).strip()
    return {"text": full_text, "paragraphs": paragraphs}


def extract_table(shape) -> Optional[Dict[str, Any]]:
    if not getattr(shape, "has_table", False):
        return None

    tbl = shape.table
    rows = []
    for r_idx, row in enumerate(tbl.rows):
        cells = []
        for c_idx, cell in enumerate(row.cells):
            cells.append(
                {
                    "col_index": c_idx,
                    "text": (cell.text or "").strip(),
                }
            )
        rows.append({"row_index": r_idx, "cells": cells})
    return {
        "n_rows": len(tbl.rows),
        "n_cols": len(tbl.columns),
        "rows": rows,
    }


def extract_image(shape) -> Optional[Dict[str, Any]]:
    """For picture shapes, extract image metadata deterministically."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None

    try:
        img = shape.image  # python-pptx Image object
        blob = img.blob
        return {
            "filename": img.filename,
            "content_type": img.content_type,
            "size_bytes": len(blob),
            "sha256": sha256_bytes(blob),
        }
    except Exception as e:
        return {"error": f"image_extract_failed: {e.__class__.__name__}: {e}"}


def extract_chart(shape) -> Optional[Dict[str, Any]]:
    """Best-effort chart extraction (not all charts/types covered)."""
    if not getattr(shape, "has_chart", False):
        return None

    try:
        ch = shape.chart
        info: Dict[str, Any] = {
            "chart_type": str(getattr(ch, "chart_type", None)),
            "has_title": bool(getattr(ch, "has_title", False)),
            "title": None,
            "series": [],
        }
        if info["has_title"] and ch.chart_title is not None:
            try:
                info["title"] = ch.chart_title.text_frame.text
            except Exception:
                info["title"] = None

        # series/category extraction varies by chart type; try common cases
        try:
            for s_idx, s in enumerate(ch.series):
                s_item: Dict[str, Any] = {"series_index": s_idx, "name": getattr(s, "name", None), "values": []}
                try:
                    for v in s.values:
                        s_item["values"].append(v)
                except Exception:
                    pass
                info["series"].append(s_item)
        except Exception:
            pass

        return info
    except Exception as e:
        return {"error": f"chart_extract_failed: {e.__class__.__name__}: {e}"}


def extract_hyperlink(shape) -> Optional[Dict[str, Any]]:
    """Hyperlink can be on shape clickAction or in text runs; we capture shape-level."""
    try:
        hl = shape.click_action.hyperlink
        if hl is None:
            return None
        if hl.address or hl.sub_address:
            return {"address": hl.address, "sub_address": hl.sub_address}
        return None
    except Exception:
        return None


def shape_bbox_points(shape) -> Dict[str, Any]:
    left = getattr(shape, "left", None)
    top = getattr(shape, "top", None)
    width = getattr(shape, "width", None)
    height = getattr(shape, "height", None)

    return {
        "left_pt": emu_to_points(left) if left is not None else None,
        "top_pt": emu_to_points(top) if top is not None else None,
        "width_pt": emu_to_points(width) if width is not None else None,
        "height_pt": emu_to_points(height) if height is not None else None,
    }


def get_shape_id(shape) -> Optional[int]:
    try:
        return int(shape.shape_id)
    except Exception:
        return None


def get_shape_name(shape) -> Optional[str]:
    try:
        return shape.name
    except Exception:
        return None


def get_shape_xml_tag(shape) -> Optional[str]:
    try:
        el: OxmlElement = shape._element  # type: ignore
        return el.tag
    except Exception:
        return None


def extract_shape_rec(shape, slide_index: int, path: str = "") -> Dict[str, Any]:
    """Recursively extract shapes; deterministic facts only."""
    shape_id = get_shape_id(shape)
    mso_type = getattr(shape, "shape_type", None)

    entry: Dict[str, Any] = {
        "slide_index": slide_index,
        "path": path,
        "shape_id": shape_id,
        "name": get_shape_name(shape),
        "shape_type": str(getattr(shape, "shape_type", None)),
        "xml_tag": get_shape_xml_tag(shape),
        "bbox": shape_bbox_points(shape),
        "rotation_deg": float(getattr(shape, "rotation", 0.0)) if getattr(shape, "rotation", None) is not None else None,
        "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
        "hyperlink": extract_hyperlink(shape),
        "text_frame": extract_textframe(shape),
        "table": extract_table(shape),
        "image": extract_image(shape),
        "chart": extract_chart(shape),
        "children": [],
    }

    if mso_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in shape.shapes:
                child_path = f"{path}/g{shape_id}" if path else f"g{shape_id}"
                entry["children"].append(extract_shape_rec(child, slide_index, path=child_path))
        except Exception as e:
            entry["children_error"] = f"group_iter_failed: {e.__class__.__name__}: {e}"

    return entry


def extract_notes_text(slide) -> Optional[str]:
    try:
        ns = slide.notes_slide
        if ns is None or ns.notes_text_frame is None:
            return None
        t = (ns.notes_text_frame.text or "").strip()
        return t if t else None
    except Exception:
        return None


def extract_presentation_ir(pptx_path: str) -> Dict[str, Any]:
    prs = Presentation(pptx_path)

    deck: Dict[str, Any] = {
        "source_file": os.path.basename(pptx_path),
        "source_path": os.path.abspath(pptx_path),
        "extracted_at_utc": datetime.utcnow().isoformat(timespec="seconds") + "Z",
        "slide_size": get_slide_size_points(prs),
        "slides": [],
    }

    for s_idx, slide in enumerate(prs.slides):
        slide_entry: Dict[str, Any] = {
            "slide_index": s_idx,
            "slide_id": None,
            "notes_text": extract_notes_text(slide),
            "shapes": [],
            "shape_count": 0,
        }

        try:
            slide_entry["slide_id"] = getattr(slide, "slide_id", None)
        except Exception:
            slide_entry["slide_id"] = None

        shapes_list = []
        for shape in slide.shapes:
            shapes_list.append(extract_shape_rec(shape, s_idx, path=""))

        slide_entry["shapes"] = shapes_list
        slide_entry["shape_count"] = len(slide.shapes)
        deck["slides"].append(slide_entry)

    return deck


def main():
    ap = argparse.ArgumentParser(description="Deterministic PPTX → JSON IR extractor")
    ap.add_argument("pptx", help="Path to input .pptx")
    ap.add_argument("-o", "--out", default=None, help="Output JSON path (default: <pptx>.ir.json)")
    ap.add_argument("--pretty", action="store_true", help="Pretty-print JSON")
    args = ap.parse_args()

    pptx_path = args.pptx
    if not os.path.isfile(pptx_path):
        raise SystemExit(f"File not found: {pptx_path}")

    out_path = args.out or (os.path.splitext(pptx_path)[0] + ".ir.json")

    ir = extract_presentation_ir(pptx_path)

    with open(out_path, "w", encoding="utf-8") as f:
        if args.pretty:
            json.dump(ir, f, ensure_ascii=False, indent=2)
        else:
            json.dump(ir, f, ensure_ascii=False)

    print(f"Wrote IR JSON to: {out_path}")
    print(f"Slides: {len(ir['slides'])}")


if __name__ == "__main__":
    main()
